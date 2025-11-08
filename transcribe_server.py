# ======================================================
# transcribe_server.py
# ======================================================
# 1) INSTALLS (pip)
# ------------------------------------------------------
# ASR + OCR deps
# pip install fastapi uvicorn[standard] httpx faster-whisper imageio-ffmpeg
# pip install pillow torch transformers accelerate                 # Florence-2-large (VLM)
# pip install paddleocr paddlepaddle                               # CPU PaddleOCR
# # For GPU PaddleOCR, install the matching wheel for your CUDA:
# #   pip install paddlepaddle-gpu==<version matching your CUDA>
#
# SINGLE-PAGE HTML extraction deps
# pip install beautifulsoup4 lxml
#
# Optional (recommended) for robust video URLs:
# pip install yt-dlp
#
# 2) HOW TO RUN
# ------------------------------------------------------
# uvicorn transcribe_server:app --host 127.0.0.1 --port 9001 --workers 1
#
# 3) ENVIRONMENT VARIABLES (optional)
# ------------------------------------------------------
#   ASR_MODEL=large-v3|medium|small|...
#   VAD_MIN_SIL_MS=500
#   KEEP_WAV=false
#   KEEP_WAV_ON_ERROR=true
#   OCR_LANGS="en,fr"
#   OCR_UPSCALE=1.5
#   OCR_RECOG_HEAD="english_g2"          # kept for compatibility (unused by Florence/Paddle)
#   OCR_ALLOWLIST="0123...ABC..."
#   VLM_MODEL_ID="microsoft/Florence-2-large"

from __future__ import annotations

# ======================================================
# IMPORTS & GLOBAL CONFIG (guarded + lazy)
# ======================================================

# --- Standard library ---
import os, math, tempfile, subprocess, io, re, unicodedata, textwrap, hashlib, json
import pathlib, shutil, sys, types, threading, datetime
from pathlib import Path
from io import BytesIO
from typing import Any, Dict, List, Optional, Tuple, Union

# --- must be at the very top, before any transformers import ---
os.environ.setdefault("HF_USE_FLASH_ATTENTION_2", "0")  # disable FA2 by default

def _install_flash_attn_stubs():
    """
    Provide importable 'flash_attn' packages with a valid __spec__/__path__ so
    importlib introspection doesn't crash on CPU-only systems.
    """
    import importlib.util as _iu

    def _stub(name: str, is_pkg: bool = False):
        m = sys.modules.get(name) or types.ModuleType(name)
        # minimal but well-formed module metadata
        m.__file__ = f"<stub {name}>"
        m.__package__ = name.rpartition(".")[0]
        if is_pkg:
            # mark as a package so submodules import cleanly
            m.__path__ = []  # type: ignore[attr-defined]
        m.__spec__ = _iu.spec_from_loader(name, loader=None)
        sys.modules[name] = m
        return m

    pkg = _stub("flash_attn", is_pkg=True)
    bert = _stub("flash_attn.bert_padding")
    ops  = _stub("flash_attn.ops")
    _stub("flash_attn.flash_attn_interface")

    # Optional no-op symbols (won't be called when we use eager attention)
    def _noimpl(*a, **k):
        raise RuntimeError("flash_attn is stubbed (CPU fallback).")
    if not hasattr(bert, "pad_input"):   setattr(bert, "pad_input", _noimpl)
    if not hasattr(bert, "unpad_input"): setattr(bert, "unpad_input", _noimpl)
    if not hasattr(ops,  "flash_attn_varlen_qkvpacked_func"):
        setattr(ops, "flash_attn_varlen_qkvpacked_func", _noimpl)

# Only now import HF after disabling flash-attn; if it still complains, add stubs then retry.
try:
    from transformers import AutoModelForCausalLM, AutoProcessor
except Exception:
    _install_flash_attn_stubs()
    from transformers import AutoModelForCausalLM, AutoProcessor

    
# --- Networking / HTTP ---
import httpx
from httpx import ConnectError

# --- FastAPI & schema ---
from fastapi import FastAPI, HTTPException, Request
from fastapi.responses import PlainTextResponse, JSONResponse
from fastapi.middleware.cors import CORSMiddleware

from pydantic import BaseModel

# --- ASR: faster-whisper + ffmpeg ---
from faster_whisper import WhisperModel
from imageio_ffmpeg import get_ffmpeg_exe

# --- Vision OCR stack: Florence-2 (VLM) + PaddleOCR fallback ---
import numpy as np
from PIL import Image
import torch

# Guard Paddle imports (optional dependency)
try:
    from paddleocr import PaddleOCR  # type: ignore
    import paddle  # type: ignore
except Exception:
    PaddleOCR = None  # type: ignore
    paddle = None

# --- HTML parsing ---
from bs4 import BeautifulSoup

# ======================================================
# FILE-TYPE / FEATURE–SPECIFIC IMPORTS (optional, graceful fallback)
# ======================================================

# --- PDF (PyMuPDF) ---
try:
    import fitz  # PyMuPDF
except Exception:
    fitz = None

# --- PPTX (python-pptx) ---
try:
    from pptx import Presentation  # type: ignore
    from pptx.enum.shapes import MSO_SHAPE_TYPE  # type: ignore
except Exception:
    Presentation = None
    MSO_SHAPE_TYPE = None

# --- DOCX (python-docx) ---
try:
    from docx import Document as Docx  # type: ignore
except Exception:
    Docx = None

# --- CSV/XLSX tables (pandas / openpyxl) ---
try:
    import pandas as pd  # type: ignore
    from tabulate import tabulate  # type: ignore
except Exception:
    pd = None  # type: ignore
    tabulate = None  # type: ignore

try:
    import openpyxl  # type: ignore
except Exception:
    openpyxl = None

# --- Text encoding detection (chardet) ---
try:
    import chardet  # type: ignore
except Exception:
    chardet = None  # type: ignore

# --- Local VLM captioning pipeline (Transformers pipeline) ---
try:
    from transformers import pipeline as hf_pipeline  # type: ignore
except Exception:
    hf_pipeline = None

# --- Legacy .doc (OLE) metadata ---
try:
    import olefile as _olefile  # pip install olefile
except Exception:
    _olefile = None

# --- RTF (striprtf optional) ---
try:
    from striprtf.striprtf import rtf_to_text  # pip install striprtf
except Exception:
    rtf_to_text = None

# --- HF snapshot (optional) ---
try:
    from huggingface_hub import snapshot_download
except Exception:
    snapshot_download = None  # type: ignore

# ------------------------------------------------------
# ASR CONFIG & DEVICE / PRECISION
# ------------------------------------------------------
MODEL_SIZE = os.getenv("ASR_MODEL", "large-v3")
VAD_MIN_SIL_MS = int(os.getenv("VAD_MIN_SIL_MS", "500"))
KEEP_WAV = os.getenv("KEEP_WAV", "false").lower() == "true"
KEEP_WAV_ON_ERROR = os.getenv("KEEP_WAV_ON_ERROR", "true").lower() == "true"

# Root for all local model snapshots (default: ./models next to this script)
MODELS_ROOT = os.getenv("MODELS_ROOT", str(Path.cwd() / "models"))
Path(MODELS_ROOT).mkdir(parents=True, exist_ok=True)

# ------------------ Florence-2 (primary VLM OCR + caption) ------------------

# Model selection
VLM_MODEL_ID = os.getenv("VLM_MODEL_ID", "microsoft/Florence-2-large")

# Optional absolute/host-mounted path to a pre-downloaded model dir (highest priority)
FLORENCE_LOCAL_DIR = os.getenv("FLORENCE_LOCAL_DIR")  # e.g., "/models/florence2"

# Single control knob
FLORENCE_ALLOW_DOWNLOAD = os.getenv("FLORENCE_ALLOW_DOWNLOAD", "").lower() == "true"
FLORENCE_REVISION = os.getenv("FLORENCE_REVISION", "main")

# Snapshot location is derived from the model id (…/models/<last-segment>)
def _default_florence_snapshot_dir(model_id: str) -> Path:
    slug = (model_id.split("/")[-1] or "florence2-large").strip().replace(" ", "_")
    return Path(MODELS_ROOT) / slug

FLORENCE_SNAPSHOT_DIR = Path(
    os.getenv("FLORENCE_SNAPSHOT_DIR", str(_default_florence_snapshot_dir(VLM_MODEL_ID)))
)

# Extreme-detailed caption tweaks
VLM_CAPTION_TASK = os.getenv("VLM_CAPTION_TASK", "<MORE_DETAILED_CAPTION>")
VLM_MAX_TOKENS_CAPTION = int(os.getenv("VLM_MAX_TOKENS_CAPTION", "768"))   # try 512–1024
VLM_REGION_CAPTION_TOPK = int(os.getenv("VLM_REGION_CAPTION_TOPK", "5"))   # caption top-N regions by area

def _parse_dtype(name: str):
    name = (name or "").lower()
    if name in ("fp16", "float16"): return torch.float16
    if name in ("bf16", "bfloat16"): return torch.bfloat16
    return torch.float32

FLORENCE_DEVICE = os.getenv("FLORENCE_DEVICE") or ("cuda:0" if torch.cuda.is_available() else "cpu")
FLORENCE_DTYPE  = _parse_dtype(os.getenv("FLORENCE_DTYPE") or ("float16" if torch.cuda.is_available() else "float32"))

# Lazy-loaded Florence singletons
_FLORENCE_LOCK = threading.Lock()
_FLORENCE_MODEL: Optional[AutoModelForCausalLM] = None
_FLORENCE_PROCESSOR: Optional[AutoProcessor] = None
_FLORENCE_LAST_ERROR: Optional[str] = None  # <-- debug surface

def _ensure_florence_path() -> Tuple[str, bool]:
    """
    Decide and prepare the source path for Florence, returning (source, is_local_path).

    Precedence:
      1) FLORENCE_LOCAL_DIR (if provided and non-empty)
      2) Snapshot dir (FLORENCE_SNAPSHOT_DIR):
           - If exists: use it
           - If missing:
               * If FLORENCE_ALLOW_DOWNLOAD=true → snapshot_download, then use it
               * If FLORENCE_ALLOW_DOWNLOAD=false → snapshot_download anyway (warn loudly), then use it
      3) If snapshot download not possible (e.g., no huggingface_hub), fall back to repo id
         (Transformers may download to HF cache).
    """
    import warnings

    def _warn(msg: str):
        try:
            warnings.warn(msg)
        finally:
            try:
                print(f"[florence2][snapshot] WARNING: {msg}", file=sys.stderr, flush=True)
            except Exception:
                pass

    # 1) Pre-mounted local dir
    if FLORENCE_LOCAL_DIR:
        p = Path(FLORENCE_LOCAL_DIR)
        if p.exists() and any(p.rglob("*")):
            return str(p), True

    # 2) Snapshot dir
    snap = FLORENCE_SNAPSHOT_DIR
    if snap.exists() and any(snap.rglob("*")):
        return str(snap), True

    # Snapshot missing → we will create/populate it
    if not FLORENCE_ALLOW_DOWNLOAD:
        _warn(
            f"No local snapshot found at '{snap}'. FLORENCE_ALLOW_DOWNLOAD=false, "
            f"but proceeding with a one-time download to populate the local cache."
        )

    if snapshot_download is not None:
        try:
            snap.parent.mkdir(parents=True, exist_ok=True)
            snapshot_download(
                repo_id=VLM_MODEL_ID,
                revision=FLORENCE_REVISION,
                local_dir=str(snap),
                local_dir_use_symlinks=False,
            )
            return str(snap), True
        except Exception as e:
            _warn(f"{type(e).__name__}: {e}. Falling back to repo id '{VLM_MODEL_ID}'.")
    else:
        _warn("huggingface_hub not installed; cannot snapshot. Falling back to repo id.")

    # 3) Fall back to repo id (Transformers may download via standard cache)
    return VLM_MODEL_ID, False

def _get_florence() -> Tuple[Optional[AutoModelForCausalLM], Optional[AutoProcessor]]:
    """
    Thread-safe lazy loader for Florence with CPU-safe attention (no flash-attn).
    - Chooses `attn_implementation` based on device (override via VLM_ATTN_IMPL).
    - If a remote modeling file tries to import flash_attn, we stub it and retry.
    """
    global _FLORENCE_MODEL, _FLORENCE_PROCESSOR, _FLORENCE_LAST_ERROR
    if _FLORENCE_MODEL is not None and _FLORENCE_PROCESSOR is not None:
        return _FLORENCE_MODEL, _FLORENCE_PROCESSOR

    # Choose an attention impl that never requires flash-attn on CPU.
    default_impl = "eager" if str(FLORENCE_DEVICE).startswith("cpu") else "sdpa"
    ATTN_IMPL = os.getenv("VLM_ATTN_IMPL", default_impl)

    with _FLORENCE_LOCK:
        if _FLORENCE_MODEL is None or _FLORENCE_PROCESSOR is None:
            try:
                source, is_local = _ensure_florence_path()

                def _load(attn_impl: str):
                    model = AutoModelForCausalLM.from_pretrained(
                        source,
                        torch_dtype=FLORENCE_DTYPE,
                        trust_remote_code=True,
                        local_files_only=bool(is_local),
                        revision=(FLORENCE_REVISION if not is_local else None),
                        attn_implementation=attn_impl,   # ← key line: avoid flash-attn
                    ).to(FLORENCE_DEVICE)
                    proc = AutoProcessor.from_pretrained(
                        source,
                        trust_remote_code=True,
                        local_files_only=bool(is_local),
                        revision=(FLORENCE_REVISION if not is_local else None),
                    )
                    return model, proc

                try:
                    _FLORENCE_MODEL, _FLORENCE_PROCESSOR = _load(ATTN_IMPL)  # uses attn_implementation
                    _FLORENCE_LAST_ERROR = None
                except Exception as e:
                    msg = f"{type(e).__name__}: {e}"
                    if "flash_attn" in msg.lower():
                        _install_flash_attn_stubs()
                        _FLORENCE_MODEL, _FLORENCE_PROCESSOR = _load("eager")
                        _FLORENCE_LAST_ERROR = None
                    else:
                        raise

            except Exception as e:
                _FLORENCE_MODEL = None
                _FLORENCE_PROCESSOR = None
                _FLORENCE_LAST_ERROR = f"{type(e).__name__}: {e}"

    return _FLORENCE_MODEL, _FLORENCE_PROCESSOR


# ------------------ PaddleOCR (fallback OCR) ------------------
_PADDLE: Optional["PaddleOCR"] = None
_PADDLE_LANG: str = (os.getenv("OCR_LANGS", "en").split(",")[0] or "en").strip()

PADDLE_LANG = _PADDLE_LANG  # alias for downstream references/logging

def _get_paddle() -> Optional["PaddleOCR"]:
    """Lazy-create a PaddleOCR instance; returns None if paddleocr isn't installed."""
    global _PADDLE
    if _PADDLE is not None:
        return _PADDLE
    if PaddleOCR is None:
        return None
    _PADDLE = PaddleOCR(use_angle_cls=True, lang=_PADDLE_LANG, show_log=False)
    return _PADDLE

def _paddle_gpu_available() -> bool:
    try:
        return bool(paddle) and paddle.device.is_compiled_with_cuda() and (paddle.device.cuda.device_count() > 0)
    except Exception:
        return False

PADDLE_USE_GPU = _paddle_gpu_available()


# ------------------ faster-whisper snapshot (ASR) ------------------

# Ensure MODELS_ROOT exists (default: ./models next to this script)
try:
    MODELS_ROOT  # type: ignore[name-defined]
except NameError:
    MODELS_ROOT = os.getenv("MODELS_ROOT", str(Path.cwd() / "models"))
    Path(MODELS_ROOT).mkdir(parents=True, exist_ok=True)

# Make sure snapshot_download is available (graceful if not installed)
try:
    snapshot_download  # type: ignore[name-defined]
except NameError:
    try:
        from huggingface_hub import snapshot_download  # pip install huggingface_hub
    except Exception:
        snapshot_download = None  # type: ignore[assignment]

# Default model SIZE → faster-whisper-large-v3
MODEL_SIZE = os.getenv("ASR_MODEL", "large-v3")  # ← default

def detect_device_and_precision():
    """Pick device & precision for faster-whisper (CUDA → float16, else CPU int8_float32)."""
    try:
        import ctranslate2
        has_cuda = ctranslate2.get_cuda_device_count() > 0
    except Exception:
        has_cuda = False
    if has_cuda:
        return "cuda", "float16"
    return "cpu", "int8_float32"

DEVICE, COMPUTE_TYPE = detect_device_and_precision()

# Single control knob (snapshot semantics)
WHISPER_ALLOW_DOWNLOAD = os.getenv("WHISPER_ALLOW_DOWNLOAD", "false").lower() == "true"
WHISPER_REVISION = os.getenv("WHISPER_REVISION", "main")

# Optional absolute/host-mounted local dir (highest priority)
WHISPER_LOCAL_DIR = os.getenv("WHISPER_LOCAL_DIR")  # e.g., "/models/faster-whisper-large-v3"

# Snapshot dir under ./models/faster-whisper-<MODEL_SIZE> by default
WHISPER_SNAPSHOT_DIR = Path(os.getenv(
    "WHISPER_SNAPSHOT_DIR",
    str(Path(MODELS_ROOT) / f"faster-whisper-{MODEL_SIZE}")
))

# Map common sizes to their repo ids (default path resolves to Systran/faster-whisper-<size>)
WHISPER_REPO_BY_SIZE = {
    "large-v3": "Systran/faster-whisper-large-v3",
    # "large-v2": "Systran/faster-whisper-large-v2",
    # "distil-large-v3": "Systran/faster-distil-whisper-large-v3",
    # "medium": "Systran/faster-whisper-medium",
    # "medium.en": "Systran/faster-whisper-medium.en",
    # "small": "Systran/faster-whisper-small",
    # "small.en": "Systran/faster-whisper-small.en",
    # "base": "Systran/faster-whisper-base",
    # "base.en": "Systran/faster-whisper-base.en",
    # "tiny": "Systran/faster-whisper-tiny",
    # "tiny.en": "Systran/faster-whisper-tiny.en",
}

_WHISPER_SNAPSHOT_LAST_ERROR: Optional[str] = None

def _ensure_whisper_path(model_size: str) -> Optional[str]:
    """
    Decide and prepare the source path for faster-whisper; return a local path if available,
    else None to let WhisperModel resolve by repo id.
    """
    import sys, warnings

    def _warn(msg: str):
        try:
            warnings.warn(msg)
        except Exception:
            pass
        try:
            print(f"[faster-whisper][snapshot] WARNING: {msg}", file=sys.stderr, flush=True)
        except Exception:
            pass

    global _WHISPER_SNAPSHOT_LAST_ERROR

    # 1) Pre-mounted local dir
    if WHISPER_LOCAL_DIR:
        p = Path(WHISPER_LOCAL_DIR)
        if p.exists() and any(p.rglob("*")):
            return str(p)

    # 2) Snapshot dir
    snap = WHISPER_SNAPSHOT_DIR
    if snap.exists() and any(snap.rglob("*")):
        return str(snap)

    # Snapshot missing → we will create/populate it
    if not WHISPER_ALLOW_DOWNLOAD:
        _warn(
            f"No local snapshot found at '{snap}'. WHISPER_ALLOW_DOWNLOAD=false, "
            f"but proceeding with a one-time download to populate the local cache."
        )

    if snapshot_download is None:
        _WHISPER_SNAPSHOT_LAST_ERROR = "huggingface_hub not installed; cannot snapshot faster-whisper."
        _warn(_WHISPER_SNAPSHOT_LAST_ERROR)
        return None

    try:
        repo_id = WHISPER_REPO_BY_SIZE.get(model_size, f"Systran/faster-whisper-{model_size}")
        snap.parent.mkdir(parents=True, exist_ok=True)
        snapshot_download(
            repo_id=repo_id,
            revision=WHISPER_REVISION,
            local_dir=str(snap),
            local_dir_use_symlinks=False,
        )
        return str(snap)
    except Exception as e:
        _WHISPER_SNAPSHOT_LAST_ERROR = f"{type(e).__name__}: {e}"
        _warn(_WHISPER_SNAPSHOT_LAST_ERROR)
        return None

# Prefer local snapshot path; fall back to the model name
try:
    _whisper_src = _ensure_whisper_path(MODEL_SIZE)
except Exception:
    _whisper_src = None  # fallback path

# --- Lazy Whisper loader (replaces global ASR_MODEL) ---
_ASR_MODEL = None

def _get_asr_model():
    global _ASR_MODEL
    if _ASR_MODEL is None:
        # Re-detect to be safe at first use (GPU vs CPU)
        device, compute_type = detect_device_and_precision()
        src = _whisper_src or MODEL_SIZE
        _ASR_MODEL = WhisperModel(src, device=device, compute_type=compute_type)
    return _ASR_MODEL

WHISPER_SNAPSHOT_ENABLED = bool(_whisper_src)  # True if we’re using a local snapshot dir

# Force default ON unless explicitly disabled via env (EAGER_SNAPSHOT_WARMUP=false)
os.environ.setdefault("EAGER_SNAPSHOT_WARMUP", "true")

# ---- EAGER SNAPSHOT WARMUP (import-time, optional) -------------------------
# Ensures local snapshot folders exist; does NOT load models into memory.
# Guarded so uvicorn --reload / multiple imports don’t run it twice.
if os.getenv("EAGER_SNAPSHOT_WARMUP", "true").lower() == "true" and not globals().get("_WARMUP_DONE"):
    globals()["_WARMUP_DONE"] = True
    try:
        _ = _ensure_whisper_path(MODEL_SIZE)
    except Exception:
        pass
    try:
        _ = _ensure_florence_path()
    except Exception:
        pass
# ---------------------------------------------------------------------------

# ===========================================================================================================

# ======================================================
# STEP 2 — HELPERS & EXTRACTORS
# ======================================================

# ------------------------------------------------------
# GENERIC HELPERS (HASH / TIME)
# ------------------------------------------------------
def _sha256_file(path: str | os.PathLike) -> str:
    """Return SHA-256 checksum of a file (streamed in 1 MiB chunks)."""
    h = hashlib.sha256()
    with open(path, "rb") as f:
        for chunk in iter(lambda: f.read(1 << 20), b""):
            h.update(chunk)
    return h.hexdigest()

def _now_iso() -> str:
    """Current timestamp in ISO-8601 (seconds precision)."""
    return datetime.datetime.now().isoformat(timespec="seconds")

# --------------------------------------------
# FILE & DOC TIMESTAMPS HELPERS
# --------------------------------------------
def _dt_to_iso(dt) -> str:
    import datetime as _dt
    try:
        if isinstance(dt, _dt.datetime):
            return dt.replace(tzinfo=None).isoformat(timespec="seconds")
        if isinstance(dt, (int, float)):
            return _dt.datetime.fromtimestamp(float(dt)).isoformat(timespec="seconds")
        if isinstance(dt, str):
            return dt.strip()
    except Exception:
        pass
    return ""

def _parse_pdf_date(s: str) -> str:
    s = (s or "").strip()
    if s.startswith("D:"):
        s = s[2:]
    try:
        y = int(s[0:4]); mo = int(s[4:6]); d = int(s[6:8])
        h = int(s[8:10]) if len(s) >= 10 else 0
        mi = int(s[10:12]) if len(s) >= 12 else 0
        se = int(s[12:14]) if len(s) >= 14 else 0
        import datetime as _dt
        return _dt.datetime(y, mo, d, h, mi, se).isoformat(timespec="seconds")
    except Exception:
        return (s or "").strip()

def _collect_doc_times(path: str, doc_type: str) -> dict:
    """
    Filesystem times for all; embedded times for: pdf, pptx, docx, xlsx, doc.
    CSV/TXT have no embedded metadata, so only filesystem times will be present.
    """
    p = pathlib.Path(path)
    st = p.stat()
    fs_created  = _dt_to_iso(st.st_ctime)
    fs_modified = _dt_to_iso(st.st_mtime)
    fs_accessed = _dt_to_iso(st.st_atime)

    embedded_created = ""
    embedded_modified = ""

    try:
        if doc_type == "pdf" and fitz is not None:
            with fitz.open(path) as _pdf:
                meta = _pdf.metadata or {}
                embedded_created  = _parse_pdf_date(meta.get("creationDate", "") or meta.get("CreationDate", ""))
                embedded_modified = _parse_pdf_date(meta.get("modDate", "") or meta.get("ModDate", ""))

        elif doc_type == "pptx" and Presentation is not None:
            try:
                prs = Presentation(path)
                cp = prs.core_properties
                embedded_created  = _dt_to_iso(getattr(cp, "created", None))
                embedded_modified = _dt_to_iso(getattr(cp, "modified", None))
            except Exception:
                pass

        elif doc_type == "docx" and Docx is not None:
            try:
                d = Docx(path)
                cp = d.core_properties
                embedded_created  = _dt_to_iso(getattr(cp, "created", None))
                embedded_modified = _dt_to_iso(getattr(cp, "modified", None))
            except Exception:
                pass

        elif doc_type == "xlsx" and openpyxl is not None:
            try:
                wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
                props = wb.properties
                embedded_created  = _dt_to_iso(getattr(props, "created", None))
                embedded_modified = _dt_to_iso(getattr(props, "modified", None))
                try:
                    wb.close()
                except Exception:
                    pass
            except Exception:
                pass

        elif doc_type == "doc":
            # Legacy binary .doc via OLE SummaryInformation (optional)
            _ole = globals().get("_olefile", None)
            if _ole is not None:
                try:
                    with _ole.OleFileIO(path) as ole:
                        meta = ole.get_metadata()
                        embedded_created  = _dt_to_iso(getattr(meta, "create_time", None))
                        embedded_modified = _dt_to_iso(getattr(meta, "last_saved_time", None))
                except Exception:
                    pass
        # else: csv/txt → no embedded times available
    except Exception:
        pass

    last_mod_iso = embedded_modified or fs_modified
    last_mod_day = last_mod_iso.split("T")[0] if last_mod_iso else ""

    days = set()
    for t in (fs_created, fs_modified, fs_accessed, embedded_created, embedded_modified):
        if t and "T" in t:
            days.add(t.split("T")[0])

    return {
        "fs": {"created": fs_created, "modified": fs_modified, "accessed": fs_accessed},
        "embedded": {"created": embedded_created, "modified": embedded_modified},
        "last_modified_day": last_mod_day,
        "modified_days": sorted(days),
    }

# --------------------------------------------
# UPSCALING & PPTX UNITS HELPERS
# --------------------------------------------
def _bytes_from_pil(img: Image.Image, fmt: str = "PNG") -> bytes:
    buf = io.BytesIO()
    img.save(buf, format=fmt)
    return buf.getvalue()

def _maybe_upscale_icon_bytes(img_bytes: bytes, *,
                              width_px: Optional[int] = None,
                              height_px: Optional[int] = None,
                              small_min_px: int = 22,
                              target_min_px: int = 64,
                              max_scale: int = 8) -> Tuple[bytes, bool, int]:
    """
    If the image looks tiny, upscale with Lanczos so OCR/VLM can see it.
    Returns (possibly_upscaled_bytes, tiny_flag, scale_factor).
    """
    try:
        im = Image.open(io.BytesIO(img_bytes)).convert("RGB")
        w, h = im.size
        if width_px and height_px:
            # Prefer caller's geometric size_HINT from layout if provided
            w = min(w, width_px)
            h = min(h, height_px)
        m = max(1, min(w, h))
        if m < small_min_px:
            scale = max(2, min(max_scale, int(math.ceil(target_min_px / max(1, m)))))
            im2 = im.resize((im.width * scale, im.height * scale), Image.LANCZOS)
            return _bytes_from_pil(im2), True, scale
        return img_bytes, False, 1
    except Exception:
        # Fail closed — return original bytes
        return img_bytes, False, 1

_EMU_PER_IN = 914400
def _emu_to_px(emu: int, dpi: int = 96) -> int:
    return int(round((emu / _EMU_PER_IN) * dpi))

# ------------------------------------------------------
# CONTENT-TYPE HELPERS
# ------------------------------------------------------
def is_image_content_type(ctype: str) -> bool:
    """True if a MIME type looks like an image/*."""
    return ctype.startswith("image/")

def is_audio_or_video_content_type(ctype: str) -> bool:
    """True if a MIME type looks like audio/*, video/*, or common media octet-streams."""
    return ctype.startswith("audio/") or ctype.startswith("video/") or ctype in (
        "application/octet-stream",
        "application/x-mpegurl",
        "application/vnd.apple.mpegurl",
    )

def is_probably_html(ctype: str, data: bytes) -> bool:
    """Detect HTML either by Content-Type or by scanning the first 256 bytes."""
    if "text/html" in ctype:
        return True
    head = data[:256].lower()
    return head.startswith(b"<!doctype html") or b"<html" in head

# ------------------------------------------------------
# WRAPPERS (OCR / LOCAL VLM)
# ------------------------------------------------------

class _OCRWrapper:
    """
    Thin wrapper around PaddleOCR to make OCR optional and pluggable.
    - engine: "paddle" to enable OCR, "none"/None to disable
    - lang:   primary OCR language (e.g., "en" or "en,fr" → picks first "en")
    - allowlist: optional character allowlist filter applied to results
    - use_global: reuse module-level singleton if available (avoids re-init cost)
                  Preferred order:
                    1) globals()['_get_paddle']() if present (lazy singleton)
                    2) globals()['PADDLE'] if language matches
                    3) per-language cache inside this class
    """

    # Cache PaddleOCR instances per language to avoid repeated heavy inits
    _CACHE: Dict[str, "PaddleOCR"] = {}
    _LOCK = threading.Lock()

    def __init__(
        self,
        engine: Optional[str] = "paddle",
        lang: str = "en",
        allowlist: Optional[str] = None,
        use_global: bool = True,
        show_log: bool = False,
    ):
        self.engine = engine if engine in {None, "paddle"} else None
        self.lang = (lang or "en").split(",")[0].strip() or "en"
        self.allowlist = allowlist  # if None, no post-filter; else filter to these chars
        self._paddle: Optional["PaddleOCR"] = None
        self._show_log = bool(show_log)

        if self.engine == "paddle":
            if PaddleOCR is None:
                # Paddle not installed; disable gracefully
                self.engine = None
            else:
                # 1) Prefer a lazy global getter if present (_get_paddle from server)
                if use_global:
                    try:
                        _getter = globals().get("_get_paddle", None)
                        if callable(_getter):
                            inst = _getter()
                            if inst is not None:
                                # Try to verify language match if globals expose it
                                glob_lang = (
                                    globals().get("_PADDLE_LANG", None)
                                    or globals().get("PADDLE_LANG", None)
                                )
                                if not glob_lang or glob_lang == self.lang:
                                    self._paddle = inst
                    except Exception:
                        self._paddle = None

                # 2) Reuse module-level singleton if language matches
                if use_global and self._paddle is None:
                    try:
                        _glob = globals().get("PADDLE", None)
                        _glob_lang = globals().get("PADDLE_LANG", None)
                        if _glob is not None and (not _glob_lang or _glob_lang == self.lang):
                            self._paddle = _glob
                    except Exception:
                        self._paddle = None

                # 3) Reuse per-lang cache or create one (thread-safe)
                if self._paddle is None:
                    with _OCRWrapper._LOCK:
                        cached = _OCRWrapper._CACHE.get(self.lang)
                        if cached is None:
                            cached = PaddleOCR(use_angle_cls=True, lang=self.lang, show_log=self._show_log)
                            _OCRWrapper._CACHE[self.lang] = cached
                        self._paddle = cached

    def _apply_allowlist(self, text: str) -> str:
        """Filter output to allowed characters if an allowlist is set."""
        if not self.allowlist:
            return text
        allowed = set(self.allowlist)
        return "".join(ch if (ch in allowed or ch.isspace()) else "" for ch in text)

    def ocr_image_bytes(
        self,
        img_bytes: bytes,
        *,
        min_confidence: float = 0.0,          # filter weak lines (0.0–1.0)
        max_lines: Optional[int] = None,      # cap output lines (None = no cap)
        join_with: str = "\n",                # customize joiner if needed
        dedup: bool = True                    # stable de-dup of identical lines
    ) -> str:
        """
        Run OCR on raw image bytes; returns one string with newlines between regions.
        Internally converts to RGB→BGR ndarray for PaddleOCR.

        Signature is backward-compatible with prior calls (extra args are keyword-only).
        """
        if self.engine is None or self._paddle is None:
            return ""
        try:
            # Prefer ndarray (BGR) input for PaddleOCR
            img = Image.open(io.BytesIO(img_bytes)).convert("RGB")
            img_np_bgr = np.array(img)[:, :, ::-1]

            out = self._paddle.ocr(img_np_bgr, cls=True) or []
            lines: List[str] = []
            for region in out:
                for line in (region or []):
                    try:
                        # line[1] = (text, score)
                        txt = (line[1][0] if line and len(line) > 1 else "").strip()
                        conf = float(line[1][1]) if line and len(line) > 1 else 1.0
                        if not txt:
                            continue
                        if conf < min_confidence:
                            continue
                        txt = self._apply_allowlist(txt)
                        if txt:
                            lines.append(txt)
                    except Exception:
                        # best-effort; skip malformed line
                        continue

            if dedup:
                seen = set()
                deduped = []
                for ln in lines:
                    if ln not in seen:
                        seen.add(ln)
                        deduped.append(ln)
                lines = deduped

            if max_lines is not None and max_lines > 0 and len(lines) > max_lines:
                lines = lines[:max_lines]

            return join_with.join(lines)
        except Exception:
            return ""

# --- replace the existing _VLMWrapper entirely ---
class _VLMWrapper:
    def __init__(self, model: Optional[str], max_tokens: int = 1024, prompt: Optional[str] = None):
        self.model = model or VLM_MODEL_ID
        self.max_tokens = max_tokens
        self.prompt = prompt or VLM_CAPTION_TASK

    @property
    def enabled(self) -> bool:
        try:
            m, p = _get_florence()
            return m is not None and p is not None
        except Exception:
            return False

    def caption_image_bytes(self, img_bytes: bytes, prompt: Optional[str] = None) -> str:
        if not self.enabled:
            return ""
        try:
            img = Image.open(io.BytesIO(img_bytes)).convert("RGB")
        except Exception:
            return ""
        return _florence_caption(img, task=(prompt or self.prompt), max_new_tokens=self.max_tokens)

# ------------------------------------------------------
# IMAGE HELPERS & XLSX→PDF
# ------------------------------------------------------
_IMG_EXTS = (".png", ".jpg", ".jpeg", ".webp", ".bmp", ".gif", ".tif", ".tiff")

def _looks_like_image_ref(s: str) -> bool:
    """Heuristically decide if a string looks like an image path/URL/data URI."""
    s = (s or "").strip().lower()
    if not s:
        return False
    if s.startswith("data:image/"):
        return True
    if s.startswith("http://") or s.startswith("https://") or s.startswith("file://"):
        return any(s.endswith(ext) for ext in _IMG_EXTS)
    return any(s.endswith(ext) for ext in _IMG_EXTS)

def _load_image_from_value(val: Optional[str], image_root: Optional[str] = None, fetch_http: bool = False) -> Optional[bytes]:
    """
    Load image bytes from a value that could be:
    - data: URI, file:// URL, local path, or http(s) URL (if fetch_http=True).
    """
    if not val or not isinstance(val, str):
        return None
    s = val.strip()
    # data URI
    if s.lower().startswith("data:image/"):
        try:
            b64 = re.sub(r"^data:image/[^;]+;base64,", "", s, flags=re.I)
            import base64 as _b64
            return _b64.b64decode(b64)
        except Exception:
            return None
    # file:// URL
    if s.lower().startswith("file://"):
        try:
            from urllib.parse import urlparse
            p = pathlib.Path(urlparse(s).path)
            if p.exists():
                return p.read_bytes()
        except Exception:
            pass
    # local relative/absolute path
    if any(s.lower().endswith(ext) for ext in _IMG_EXTS):
        p = pathlib.Path(image_root or ".") / s
        if p.exists():
            try:
                return p.read_bytes()
            except Exception:
                pass
    # http(s) URL
    if (s.lower().startswith("http://") or s.lower().startswith("https://")) and fetch_http:
        try:
            import urllib.request as urlreq
            with urlreq.urlopen(s, timeout=10) as r:
                return r.read()
        except Exception:
            return None
    return None

def _analyze_image(img_bytes: bytes, ocr: _OCRWrapper | None, vlm: _VLMWrapper | None, prompt: Optional[str] = None) -> Dict[str, str]:
    """Run OCR and/or VLM captioning on image bytes; return dict with 'ocr' and 'caption'."""
    ocr_text = ""
    cap_text = ""
    if ocr and getattr(ocr, "engine", None):
        try:
            ocr_text = ocr.ocr_image_bytes(img_bytes) or ""
        except Exception:
            ocr_text = ""
    if vlm and getattr(vlm, "enabled", False):
        try:
            cap_text = vlm.caption_image_bytes(img_bytes, prompt=prompt) or ""
        except Exception:
            cap_text = ""
    return {"ocr": ocr_text.strip(), "caption": cap_text.strip()}

def _export_xlsx_to_pdf(xlsx_path: str, outdir: Optional[str] = None, soffice_bin: Optional[str] = None) -> Optional[str]:
    """Convert an .xlsx file to PDF via LibreOffice/soffice; return output PDF path or None."""
    bin_path = soffice_bin or shutil.which("soffice") or shutil.which("libreoffice")
    if not bin_path:
        return None
    p = pathlib.Path(xlsx_path)
    outdir = outdir or str(p.parent)
    try:
        subprocess.run([bin_path, "--headless", "--convert-to", "pdf", "--outdir", outdir, str(p)], check=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE)
    except Exception:
        return None
    pdf_guess = pathlib.Path(outdir) / (p.stem + ".pdf")
    return str(pdf_guess) if pdf_guess.exists() else None

# ------------------------------------------------------
# DOC CONVERTERS — Legacy .doc → .docx / .pdf 
# ------------------------------------------------------
def _is_rtf_bytes(b: bytes) -> bool:
    """True if bytes look like RTF ('{\\rtf' after optional BOM/whitespace)."""
    head = b.lstrip()[:8]
    return head.startswith(b'{\\rtf') or head.startswith(b'{\\RTF')

def _rtf_strip_text(rtf_bytes: bytes) -> str:
    """Best-effort RTF → plain text. Prefer striprtf; fallback to naive scrub."""
    txt = ""
    if rtf_to_text is not None:
        try:
            # RTF is ANSI; latin-1 keeps bytes round-trippable
            txt = rtf_to_text(rtf_bytes.decode("latin-1", "ignore")).strip()
            if txt:
                return txt
        except Exception:
            pass
    # Naive fallback: remove RTF control words/groups and decode \'hh escapes
    s = rtf_bytes.decode("latin-1", "ignore")
    # decode \'hh → byte
    s = re.sub(r"\\'([0-9a-fA-F]{2})", lambda m: bytes.fromhex(m.group(1)).decode("latin-1", "ignore"), s)
    # drop control words like \b, \par, \fs24, etc.
    s = re.sub(r"\\[a-zA-Z]+-?\d* ?", "", s)
    # remove braces
    s = s.replace("{", "").replace("}", "")
    # collapse whitespace
    s = re.sub(r"[ \t]+", " ", s)
    s = re.sub(r"\n{3,}", "\n\n", s)
    return s.strip()

def _rtf_iter_pict_hex(rtf_bytes: bytes):
    """
    Yield (blip, hex_string) for each {\pict ... hex ...} block.
    blip is 'png'/'jpeg'/'jpg'/'bmp' etc (best-effort from control word).
    """
    s = rtf_bytes.decode("latin-1", "ignore")
    i = 0
    n = len(s)
    while True:
        j = s.find("{\\pict", i)
        if j == -1:
            return
        # Walk to matching '}' (brace depth)
        depth = 0
        k = j
        while k < n:
            if s[k] == "{":
                depth += 1
            elif s[k] == "}":
                depth -= 1
                if depth == 0:
                    break
            k += 1
        if k >= n:
            return
        pict_block = s[j:k+1]
        # detect blip type
        blip = "bin"
        if "\\pngblip" in pict_block: blip = "png"
        elif "\\jpegblip" in pict_block: blip = "jpeg"
        elif "\\jpgblip" in pict_block: blip = "jpg"
        elif "\\dibitmap" in pict_block or "\\wmetafile" in pict_block: blip = "bmp"
        # hex payload is everything that looks like hex after last control word
        # strip control sequences
        body = re.sub(r"\\[a-zA-Z]+-?\d* ?", " ", pict_block)
        # keep only hex and whitespace
        hex_only = re.sub(r"[^0-9A-Fa-f]", " ", body)
        hex_only = re.sub(r"\s+", "", hex_only)
        if len(hex_only) >= 2 and len(hex_only) % 2 == 0:
            yield (blip, hex_only)
        i = k + 1

def _rtf_extract_images(rtf_bytes: bytes) -> List[bytes]:
    """Convert all \pict hex blobs to raw image bytes."""
    imgs: List[bytes] = []
    for _blip, hx in _rtf_iter_pict_hex(rtf_bytes):
        try:
            imgs.append(bytes.fromhex(hx))
        except Exception:
            continue
    return imgs

def _export_doc_to_docx(doc_path: str, outdir: Optional[str] = None, soffice_bin: Optional[str] = None) -> Optional[str]:
    """Convert legacy .doc → .docx via LibreOffice; return output path or None."""
    bin_path = soffice_bin or shutil.which("soffice") or shutil.which("libreoffice")
    if not bin_path:
        return None
    p = pathlib.Path(doc_path)
    outdir = outdir or str(p.parent)
    try:
        # LibreOffice auto-detects filter for .doc → .docx
        subprocess.run(
            [bin_path, "--headless", "--convert-to", "docx", "--outdir", outdir, str(p)],
            check=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE
        )
    except Exception:
        return None
    out = pathlib.Path(outdir) / (p.stem + ".docx")
    return str(out) if out.exists() else None


def _export_doc_to_pdf(doc_path: str, outdir: Optional[str] = None, soffice_bin: Optional[str] = None) -> Optional[str]:
    """Convert legacy .doc → .pdf via LibreOffice; return output path or None."""
    bin_path = soffice_bin or shutil.which("soffice") or shutil.which("libreoffice")
    if not bin_path:
        return None
    p = pathlib.Path(doc_path)
    outdir = outdir or str(p.parent)
    try:
        subprocess.run(
            [bin_path, "--headless", "--convert-to", "pdf", "--outdir", outdir, str(p)],
            check=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE
        )
    except Exception:
        return None
    out = pathlib.Path(outdir) / (p.stem + ".pdf")
    return str(out) if out.exists() else None


def _export_pptx_to_pdf(pptx_path: str, outdir: Optional[str] = None, soffice_bin: Optional[str] = None) -> Optional[str]:
    bin_path = soffice_bin or shutil.which("soffice") or shutil.which("libreoffice")
    if not bin_path:
        return None
    p = pathlib.Path(pptx_path)
    outdir = outdir or str(p.parent)
    try:
        subprocess.run([bin_path, "--headless", "--convert-to", "pdf", "--outdir", outdir, str(p)],
                       check=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE)
    except Exception:
        return None
    pdf_guess = pathlib.Path(outdir) / (p.stem + ".pdf")
    return str(pdf_guess) if pdf_guess.exists() else None


def _export_docx_to_pdf(docx_path: str, outdir: Optional[str] = None, soffice_bin: Optional[str] = None) -> Optional[str]:
    bin_path = soffice_bin or shutil.which("soffice") or shutil.which("libreoffice")
    if not bin_path:
        return None
    p = pathlib.Path(docx_path)
    outdir = outdir or str(p.parent)
    try:
        subprocess.run([bin_path, "--headless", "--convert-to", "pdf", "--outdir", outdir, str(p)],
                       check=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE)
    except Exception:
        return None
    pdf_guess = pathlib.Path(outdir) / (p.stem + ".pdf")
    return str(pdf_guess) if pdf_guess.exists() else None

# ------------------------------------------------------
# MARKDOWN RENDERING HELPERS
# ------------------------------------------------------
def _json_to_yaml(obj: Any, indent: int = 0) -> str:
    """Render a Python dict/list as a simple YAML-like string (no external deps)."""
    sp = "  " * indent
    if isinstance(obj, dict):
        out = []
        for k, v in obj.items():
            if isinstance(v, (dict, list)):
                out.append(f"{sp}{k}:")
                out.append(_json_to_yaml(v, indent + 1))
            else:
                if isinstance(v, str):
                    out.append(f'{sp}{k}: "{v}"')
                else:
                    out.append(f"{sp}{k}: {v}")
        return "\n".join(out)
    if isinstance(obj, list):
        out = []
        for v in obj:
            if isinstance(v, (dict, list)):
                out.append(f"{sp}-")
                out.append(_json_to_yaml(v, indent + 1))
            else:
                if isinstance(v, str):
                    out.append(f'{sp}- "{v}"')
                else:
                    out.append(f"{sp}- {v}")
        return "\n".join(out)
    return f"{sp}{obj}"

def _df_to_markdown(df, max_rows: int | None = None) -> str:
    """Pretty-print a DataFrame as Markdown. If max_rows<=0 or None, show ALL rows."""
    if max_rows is None or max_rows <= 0:
        if tabulate is None:
            return df.to_csv(index=False)
        return tabulate(df, headers="keys", tablefmt="pipe", showindex=False)

    if tabulate is None:
        head = df.head(max_rows)
        return head.to_csv(index=False)

    total = len(df)
    head = df.head(max_rows)
    note = f"\n\n_… {total - max_rows} more rows omitted …_" if total > max_rows else ""
    md = tabulate(head, headers="keys", tablefmt="pipe", showindex=False)
    return md + note

def _detect_group_key(records: List[Dict[str, Any]]) -> Optional[str]:
    """Pick a grouping key among page/slide/sheet if present in records."""
    for key in ("page", "slide", "sheet"):
        if any(key in r and r[key] not in (None, 0, "") for r in records):
            return key
    return None

def _sorted_group_keys(records: List[Dict[str, Any]], key: str) -> List[Any]:
    """Return sorted unique values for the given grouping key, robust to mixed types."""
    keys = {r.get(key) for r in records if r.get(key) not in (None, "")}
    try:
        return sorted(keys)
    except Exception:
        return sorted(keys, key=lambda x: str(x))

def _join_text_blocks(blocks: List[Dict[str, Any]]) -> str:
    """Join the 'text' fields from a list of record blocks, skipping empties."""
    if not blocks:
        return "_(none)_"
    parts: List[str] = []
    for b in blocks:
        txt = (b.get("text") or "").strip()
        if not txt:
            continue
        parts.append(txt)
    return "\n\n".join(parts)

def _render_markdown_note(
    title: str,
    src_path: str,
    url: Optional[str],
    records: List[Dict[str, Any]],
    mime: str,
    doc_type: str,
    languages: List[str],
    text_extractor: str,
    ocr_engine: Optional[str],
    captioner: Optional[str],
    extra_meta: Optional[Dict[str, Any]] = None,
) -> str:
    """Compose a full Markdown report with front matter, summary, and grouped contents."""
    p = pathlib.Path(src_path)
    checksum = _sha256_file(src_path)
    pages = max([r.get("page") or 0 for r in records] or [0])
    slides = max([r.get("slide") or 0 for r in records] or [0])
    sheets = len({r.get("sheet") for r in records if r.get("sheet") is not None})

    fm = {
        "title": title,
        "source": {"file": p.name, "url": url or "", "mime": mime},
        "checksum_sha256": checksum,
        "doc_type": doc_type,
        "pages": pages,
        "slides": slides,
        "sheets": sheets,
        "languages": languages,
        "generated": _now_iso(),
        "pipeline": {
            "text_extractor": text_extractor,
            "ocr_engine": ocr_engine or "",
            "captioner": captioner or "",
        },
    }
    if extra_meta:
        fm.update(extra_meta)

    lines = [
        "---",
        _json_to_yaml(fm),
        "---",
        "",
        "# Summary",
        f"- File: **{p.name}**",
        f"- Generated: **{fm['generated']}**",
        f"- Pages/Slides/Sheets: **{pages}/{slides}/{sheets}**",
        "- Notes: Native text + OCR text merged; captions added for figures where available.",
        "",
        "# Contents",
    ]

    group_key = _detect_group_key(records)
    if group_key:
        for k in _sorted_group_keys(records, group_key):
            lines.append(f"## {group_key.capitalize()} {k}")
            native = [r for r in records if r.get(group_key) == k and r.get("source") == "native"]
            lines.append("### Native")
            lines.append(_join_text_blocks(native))
            lines.append("")
            ocr = [r for r in records if r.get(group_key) == k and r.get("source") == "ocr"]
            lines.append("### OCR")
            lines.append(_join_text_blocks(ocr))
            lines.append("")
            caps = [r for r in records if r.get(group_key) == k and r.get("source") == "vlm_alt"]
            lines.append("### Figures (captions)")
            if caps:
                for r in caps:
                    t = (r.get("text","") or "").strip()
                    if t:
                        lines.append(f"- {t}")
            else:
                lines.append("_(none)_")
            lines.append("")
        # NEW: show items without the group key
        ung = [r for r in records if not r.get(group_key)]
        if ung:
            lines.append("## Other content")
            native = [r for r in ung if r.get("source") == "native"]
            ocr    = [r for r in ung if r.get("source") == "ocr"]
            caps   = [r for r in ung if r.get("source") == "vlm_alt"]
            lines.append("### Native")
            lines.append(_join_text_blocks(native)); lines.append("")
            lines.append("### OCR")
            lines.append(_join_text_blocks(ocr)); lines.append("")
            lines.append("### Figures (captions)")
            if caps:
                for r in caps:
                    t = (r.get("text","") or "").strip()
                    if t:
                        lines.append(f"- {t}")
            else:
                lines.append("_(none)_")
            lines.append("")
    else:
        lines.append("## Content")
        lines.append(_join_text_blocks(records))
    return "\n".join(lines)  

# ------------------------------------------------------
# EXTRACTORS — FILE TYPE CONSTANTS
# ------------------------------------------------------
_DOC_TYPES = {
    ".pdf": ("application/pdf", "pdf"),
    ".pptx": ("application/vnd.openxmlformats-officedocument.presentationml.presentation", "pptx"),
    ".docx": ("application/vnd.openxmlformats-officedocument.wordprocessingml.document", "docx"),
    ".doc":  ("application/msword", "doc"),  # Legacy OLE .doc
    ".txt": ("text/plain", "txt"),
    ".csv": ("text/csv", "csv"),
    ".xlsx": ("application/vnd.openxmlformats-officedocument.spreadsheetml.sheet", "xlsx"),
}

# ------------------------------------------------------
# EXTRACTORS — PDF
# ------------------------------------------------------

def _extract_pdf(
    path: str,
    ocr: _OCRWrapper,
    vlm: Optional[_VLMWrapper] = None,
    *,
    dpi: int = 300,
    char_threshold: int = 50,
    caption_pages: bool = False,
    vlm_dpi: int = 256,
    region_dpi: int = 384,
    image_limit_per_page: int = 256,   # ↑ allow many small icons
    min_region_wh: int = 22,           # now used only to pick DPI, not to skip
    small_region_dpi: int = 640,       # ↑ render tiny icons sharper
    small_region_pad: int = 2          # add a few pt padding around tiny icons
) -> List[Dict[str, Any]]:
    """
    Per page:
      • Native text
      • Page OCR (if native text sparse)
      • Image-region OCR (ALL images, tiny included)
      • Page caption (optional)
      • Image-region captions (if VLM enabled)

    NOTE: `min_region_wh` is a *soft* threshold used to switch to high-DPI rendering
          for tiny icons; regions are never skipped because they're small.
    """
    if fitz is None:
        raise RuntimeError("PyMuPDF is required for PDF support. pip install pymupdf")

    recs: List[Dict[str, Any]] = []
    doc = fitz.open(path)

    for page_idx, page in enumerate(doc, start=1):
        # -------- 1) Native text --------
        try:
            blocks = page.get_text("blocks")
            texts = [
                b[4] for b in blocks
                if isinstance(b, (list, tuple)) and len(b) >= 5 and isinstance(b[4], str) and b[4].strip()
            ]
            native = "\n".join(texts)
            if native.strip():
                recs.append({"page": page_idx, "source": "native", "section": "paragraph", "text": native})
            char_count = sum(len(t) for t in texts)
        except Exception:
            char_count = 0

        # -------- 2) Page OCR when text is sparse --------
        if ocr.engine is not None and char_count < char_threshold:
            try:
                pix = page.get_pixmap(dpi=dpi, alpha=False)
                ocr_text = ocr.ocr_image_bytes(pix.tobytes("png"))
                if ocr_text.strip():
                    recs.append({"page": page_idx, "source": "ocr", "section": "page_ocr", "text": ocr_text})
            except Exception:
                pass

        # -------- 3) Page caption (optional) --------
        if caption_pages and vlm and vlm.enabled:
            try:
                pix = page.get_pixmap(dpi=vlm_dpi, alpha=False)
                cap = vlm.caption_image_bytes(pix.tobytes("png"))
                if cap and cap.strip():
                    recs.append({"page": page_idx, "source": "vlm_alt", "section": "page_caption", "text": cap})
            except Exception:
                pass

        # -------- 4) Image-region OCR + captions (include tiny icons) --------
        try:
            raw = page.get_text("rawdict") or {}
            blocks_rd = raw.get("blocks", [])
            seen = 0

            for b in blocks_rd:
                if seen >= image_limit_per_page:
                    break
                if not isinstance(b, dict) or b.get("type") != 1:
                    continue

                bbox = b.get("bbox")
                if not bbox or len(bbox) != 4:
                    continue
                x1, y1, x2, y2 = bbox
                w = max(0.0, (x2 - x1))
                h = max(0.0, (y2 - y1))
                if w <= 0 or h <= 0:
                    continue

                # Treat tiny icons specially (no skipping; just render sharper + small padding)
                tiny = (w < min_region_wh) or (h < min_region_wh)
                use_dpi = small_region_dpi if tiny else region_dpi
                pad = small_region_pad if tiny else 0

                # Clip with padding (clamped to page)
                try:
                    clip = fitz.Rect(max(0, x1 - pad), max(0, y1 - pad), x2 + pad, y2 + pad)
                    mat = fitz.Matrix(use_dpi / 72.0, use_dpi / 72.0)
                    pix = page.get_pixmap(matrix=mat, clip=clip, alpha=False)
                    img_bytes = pix.tobytes("png")
                except Exception:
                    continue

                res = _analyze_image(img_bytes, ocr=ocr, vlm=vlm, prompt=None)

                meta = {
                    "page": page_idx,
                    "bbox": [int(x1), int(y1), int(x2), int(y2)],
                    "tiny_icon": bool(tiny),
                    "dpi_used": int(use_dpi),
                }
                if res.get("ocr"):
                    recs.append({**meta, "source": "ocr", "section": "image_region", "text": res["ocr"]})
                if res.get("caption"):
                    recs.append({**meta, "source": "vlm_alt", "section": "image_region_caption", "text": res["caption"]})

                seen += 1

        except Exception:
            pass

    return recs


def _caption_pdf_pages(path: str, vlm: _VLMWrapper, dpi: int = 256) -> List[Dict[str, Any]]:
    """Optional: caption each page image with a local VLM."""
    if fitz is None or not vlm or not vlm.enabled:
        return []
    recs: List[Dict[str, Any]] = []
    try:
        doc = fitz.open(path)
    except Exception:
        return recs
    for i, page in enumerate(doc, start=1):
        try:
            pix = page.get_pixmap(dpi=dpi)
            cap = vlm.caption_image_bytes(pix.tobytes("png"))
            if cap and cap.strip():
                recs.append({"page": i, "source": "vlm_alt", "section": "page_caption", "text": cap})
        except Exception:
            continue
    return recs

# ------------------------------------------------------
# EXTRACTORS — PPTX
# ------------------------------------------------------
def _extract_pptx(
    path: str,
    ocr: _OCRWrapper,
    vlm: Optional[_VLMWrapper] = None,
    *,
    tiny_min_px: int = 22,
    target_min_px: int = 64,
    max_icons_per_slide: int = 512
) -> List[Dict[str, Any]]:
    """Extract text/notes; OCR+caption ALL images incl. tiny icons (group-aware)."""
    if Presentation is None:
        raise RuntimeError("python-pptx is required for PPTX support. pip install python-pptx")

    def _iter_shapes(container):
        for sh in getattr(container, "shapes", []):
            try:
                if MSO_SHAPE_TYPE and getattr(sh, "shape_type", None) == MSO_SHAPE_TYPE.GROUP:
                    # Recurse into group
                    for g in _iter_shapes(sh):
                        yield g
                else:
                    yield sh
            except Exception:
                # On any weird shape, still yield it (best effort)
                yield sh

    pres = Presentation(path)
    recs: List[Dict[str, Any]] = []

    for slide_idx, slide in enumerate(pres.slides, start=1):
        # Notes
        try:
            if slide.has_notes_slide and slide.notes_slide and slide.notes_slide.notes_text_frame:
                nt = slide.notes_slide.notes_text_frame.text
                if nt and nt.strip():
                    recs.append({"slide": slide_idx, "source": "native", "section": "notes", "text": nt})
        except Exception:
            pass

        # Text & images (group-aware)
        seen_imgs = 0
        for s in _iter_shapes(slide):
            try:
                # Text from shapes
                if getattr(s, "has_text_frame", False) and s.text_frame and s.text_frame.text and s.text_frame.text.strip():
                    recs.append({"slide": slide_idx, "source": "native", "section": "shape_text", "text": s.text_frame.text})

                # Pictures
                if hasattr(s, "image") and getattr(s, "image", None) is not None:
                    if seen_imgs >= max_icons_per_slide:
                        continue
                    img_bytes = s.image.blob
                    # Layout hint for tiny detection (EMU -> px @96dpi)
                    w_px = _emu_to_px(getattr(s, "width", 0) or 0, dpi=96)
                    h_px = _emu_to_px(getattr(s, "height", 0) or 0, dpi=96)
                    up_bytes, tiny, scale = _maybe_upscale_icon_bytes(
                        img_bytes,
                        width_px=w_px or None,
                        height_px=h_px or None,
                        small_min_px=tiny_min_px,
                        target_min_px=target_min_px,
                    )
                    alt = ""
                    try:
                        alt = (getattr(s, "alternative_text", "") or "").strip()
                    except Exception:
                        pass

                    res = _analyze_image(up_bytes, ocr=ocr, vlm=vlm, prompt=None)
                    meta = {
                        "slide": slide_idx,
                        "source_bbox_px": [w_px, h_px],
                        "tiny_icon": bool(tiny),
                        "upscale_factor": int(scale),
                    }
                    if alt:
                        meta["alt"] = alt
                    if res.get("ocr"):
                        recs.append({**meta, "source": "ocr", "section": "image", "text": res["ocr"]})
                    if res.get("caption"):
                        recs.append({**meta, "source": "vlm_alt", "section": "image_caption", "text": res["caption"]})
                    seen_imgs += 1
            except Exception:
                continue

    return recs

# ------------------------------------------------------
# EXTRACTORS — DOC (legacy .doc)
# ------------------------------------------------------
def _extract_doc_legacy(
    path: str,
    ocr: _OCRWrapper,
    vlm: Optional[_VLMWrapper] = None,
    *,
    soffice_bin: Optional[str] = None,
    dpi: int = 300,
    ocr_threshold: int = 50,
    caption_pages: bool = False,
    vlm_dpi: int = 256,
) -> List[Dict[str, Any]]:
    """
    Best-effort extractor for legacy .doc:

      0) **NEW**: If file is actually RTF, parse text + embedded \\pict images (no external tools)
      1) Try LibreOffice convert → .docx, then reuse _extract_docx (captures text + images)
      2) Else try LibreOffice convert → .pdf, then reuse _extract_pdf (text/OCR + image regions)
      3) Else fallback to CLI text extractors (antiword/catdoc/wvText) — text only
    """
    recs: List[Dict[str, Any]] = []

    # 0) RTF disguised as .doc → parse in-process
    try:
        raw = pathlib.Path(path).read_bytes()
        if _is_rtf_bytes(raw):
            # Text
            txt = _rtf_strip_text(raw)
            if txt:
                recs.append({"source": "native", "section": "paragraph", "text": txt})
            # Images (OCR/VLM)
            imgs = _rtf_extract_images(raw)
            for b in imgs:
                try:
                    # If tiny icons (often 16–24 px), upscale a bit for OCR/VLM
                    try:
                        pil = Image.open(io.BytesIO(b)).convert("RGB")
                        w, h = pil.size
                        if min(w, h) < 22:
                            pil = pil.resize((max(22, w*4), max(22, h*4)), Image.LANCZOS)
                            b = _bytes_from_pil(pil, "PNG")
                    except Exception:
                        pass
                    res = _analyze_image(b, ocr=ocr, vlm=vlm, prompt=None)
                    if res.get("ocr"):
                        recs.append({"source": "ocr", "section": "image", "text": res["ocr"]})
                    if res.get("caption"):
                        recs.append({"source": "vlm_alt", "section": "image_caption", "text": res["caption"]})
                except Exception:
                    continue
            # If we got anything at all, return here.
            if recs:
                return recs
    except Exception:
        # fall through to converters
        pass

    # 1) Prefer .docx conversion (retains structure; image sweep works)
    docx_path = _export_doc_to_docx(path, soffice_bin=soffice_bin)
    if docx_path and os.path.exists(docx_path):
        try:
            return _extract_docx(
                docx_path,
                ocr=ocr,
                vlm=vlm,
                tiny_min_px=22,
                target_min_px=64,
                max_images=1024,
            )
        except Exception:
            pass  # fall through

    # 2) PDF route (good for OCR + page/image-region passes)
    pdf_path = _export_doc_to_pdf(path, soffice_bin=soffice_bin)
    if pdf_path and os.path.exists(pdf_path):
        try:
            return _extract_pdf(
                pdf_path,
                ocr=ocr,
                vlm=vlm,
                dpi=dpi,
                char_threshold=ocr_threshold,
                caption_pages=caption_pages,
                vlm_dpi=vlm_dpi,
                region_dpi=384,
                image_limit_per_page=256,
                min_region_wh=22,
                small_region_dpi=640,
                small_region_pad=2,
            )
        except Exception:
            pass  # fall through

    # 3) Last-resort: plain text via system CLIs (no images)
    def _run(cmd: List[str]) -> Optional[str]:
        try:
            p = subprocess.run(cmd, stdout=subprocess.PIPE, stderr=subprocess.PIPE, check=True)
            out = p.stdout.decode("utf-8", "ignore").strip()
            return out or None
        except Exception:
            return None

    txt = None
    if shutil.which("antiword"):
        txt = _run(["antiword", "-w", "0", path]) or txt
    if txt is None and shutil.which("catdoc"):
        txt = _run(["catdoc", "-w", path]) or txt
    if txt is None and shutil.which("wvText"):
        with tempfile.TemporaryDirectory() as td:
            out_txt = pathlib.Path(td) / "out.txt"
            _ = _run(["wvText", path, str(out_txt)])
            try:
                if out_txt.exists():
                    txt = out_txt.read_text("utf-8", errors="ignore").strip() or None
            except Exception:
                pass

    if txt:
        recs.append({"source": "native", "section": "paragraph", "text": txt})
    else:
        recs.append({
            "source": "native",
            "section": "paragraph",
            "text": "(Could not extract content from legacy .doc and no converters/CLIs were available.)"
        })
    return recs

# ------------------------------------------------------
# EXTRACTORS — DOCX
# ------------------------------------------------------
def _extract_docx(
    path: str,
    ocr: _OCRWrapper,
    vlm: Optional[_VLMWrapper] = None,
    *,
    tiny_min_px: int = 22,
    target_min_px: int = 64,
    max_images: int = 1024
) -> List[Dict[str, Any]]:
    """
    Extract paragraphs/tables; OCR+caption images from:
      • main document body (doc.part.rels)
      • headers/footers across sections
    Tiny icons are upscaled before OCR/VLM.
    """
    if Docx is None:
        raise RuntimeError("python-docx is required for DOCX support. pip install python-docx")

    doc = Docx(path)
    recs: List[Dict[str, Any]] = []

    # 1) Native text (paragraphs + simple table rows)
    for p in doc.paragraphs:
        try:
            if p.text and p.text.strip():
                recs.append({"source": "native", "section": "paragraph", "text": p.text})
        except Exception:
            continue
    for tbl in doc.tables:
        for r in tbl.rows:
            try:
                row_txt = " | ".join((c.text or "").strip() for c in r.cells)
                if row_txt.strip():
                    recs.append({"source": "native", "section": "table", "text": row_txt})
            except Exception:
                continue

    # Helper to sweep a docx part (body/header/footer) for images
    def _sweep_part_images(part, where_label: str, limit_left: List[int]):
        # limit_left is a single-item list so inner fn can modify it (closure)
        if not part or not hasattr(part, "rels"):
            return
        for rel in part.rels.values():
            if limit_left[0] <= 0:
                return
            try:
                if "image" not in getattr(rel, "reltype", ""):
                    continue
                img_part = rel.target_part
                img_bytes = img_part.blob
                # We seldom know intended layout size in DOCX reliably; use intrinsic size
                try:
                    pil = Image.open(io.BytesIO(img_bytes)).convert("RGB")
                    w_px, h_px = pil.size
                except Exception:
                    w_px = h_px = 0

                up_bytes, tiny, scale = _maybe_upscale_icon_bytes(
                    img_bytes,
                    width_px=w_px or None,
                    height_px=h_px or None,
                    small_min_px=tiny_min_px,
                    target_min_px=target_min_px,
                )
                res = _analyze_image(up_bytes, ocr=ocr, vlm=vlm, prompt=None)

                meta = {
                    "source": where_label,
                    "section": "image",
                    "location": where_label,
                    "tiny_icon": bool(tiny),
                    "upscale_factor": int(scale),
                    "source_bbox_px": [w_px, h_px],
                }
                if res.get("ocr"):
                    recs.append({**meta, "source": "ocr", "section": "image", "text": res["ocr"]})
                if res.get("caption"):
                    recs.append({**meta, "source": "vlm_alt", "section": "image_caption", "text": res["caption"]})

                limit_left[0] -= 1
            except Exception:
                continue

    # 2) Images in main body
    remaining = [max_images]
    _sweep_part_images(doc.part, "body", remaining)

    # 3) Images in headers/footers (all sections)
    try:
        for sec in doc.sections:
            if remaining[0] <= 0:
                break
            try:
                _sweep_part_images(sec.header.part, "header", remaining)
            except Exception:
                pass
            try:
                _sweep_part_images(sec.footer.part, "footer", remaining)
            except Exception:
                pass
    except Exception:
        pass

    return recs

# ------------------------------------------------------
# EXTRACTORS — TXT
# ------------------------------------------------------
def _extract_txt(path: str, encoding: Optional[str] = None) -> List[Dict[str, Any]]:
    """Load plain-text with charset detection fallback (chardet if available)."""
    if encoding is None:
        if chardet is None:
            enc = "utf-8"
        else:
            with open(path, "rb") as f:
                enc = chardet.detect(f.read(4096)).get("encoding") or "utf-8"
    else:
        enc = encoding
    with open(path, "r", encoding=enc, errors="replace") as f:
        content = f.read()
    return [{"source": "native", "section": "paragraph", "text": content}]

# ------------------------------------------------------
# EXTRACTORS — CSV (ENHANCED: linked images OCR/VLM)
# ------------------------------------------------------
def _extract_csv_enhanced(path: str, ocr: _OCRWrapper | None, vlm: Optional[_VLMWrapper], encoding: Optional[str] = None, max_rows: int = 50, image_columns: Optional[str] = None, image_root: Optional[str] = None, fetch_http: bool = False) -> List[Dict[str, Any]]:
    """Read CSV to Markdown table; detect/link image columns for OCR and captions."""
    if pd is None:
        raise RuntimeError("pandas required for CSV support. pip install pandas")
    if encoding is None:
        if chardet is None:
            enc = "utf-8"
        else:
            with open(path, "rb") as f:
                enc = chardet.detect(f.read(4096)).get("encoding") or "utf-8"
    else:
        enc = encoding
    df = pd.read_csv(path, encoding=enc)
    recs: List[Dict[str, Any]] = []
    md_table = _df_to_markdown(df, max_rows=max_rows)
    recs.append({"sheet": pathlib.Path(path).name, "source": "native", "section": "table", "text": md_table})
    cols: List[str] = []
    if image_columns:
        cols = [c.strip() for c in image_columns.split(",") if c.strip() in df.columns]
    else:
        for c in df.columns:
            try:
                s = df[c].dropna().astype(str).head(100)
                if any(_looks_like_image_ref(v) for v in s):
                    cols.append(c)
            except Exception:
                continue
    if cols:
        for idx, row in df.iterrows():
            for c in cols:
                val = row.get(c, None)
                if val is None:
                    continue
                sval = str(val)
                if not _looks_like_image_ref(sval):
                    continue
                b = _load_image_from_value(sval, image_root=image_root, fetch_http=fetch_http)
                if not b:
                    continue
                res = _analyze_image(b, ocr=ocr, vlm=vlm, prompt=None)
                if res.get("ocr"):
                    recs.append({"sheet": pathlib.Path(path).name, "source": "ocr", "section": "linked_image", "text": f"[row {idx+1} col {c}] {res['ocr']}"})
                if res.get("caption"):
                    recs.append({"sheet": pathlib.Path(path).name, "source": "vlm_alt", "section": "linked_image_caption", "text": f"[row {idx+1} col {c}] {res['caption']}"})
    return recs

# ------------------------------------------------------
# EXTRACTORS — XLSX 
# ------------------------------------------------------
def _extract_xlsx_enhanced(path: str, ocr: _OCRWrapper | None, vlm: Optional[_VLMWrapper], max_rows: int = 50, image_root: Optional[str] = None, fetch_http: bool = False) -> List[Dict[str, Any]]:
    """Read each sheet to Markdown; OCR/caption embedded or linked images if present."""
    if pd is None:
        raise RuntimeError("pandas required for XLSX support. pip install pandas openpyxl")
    recs: List[Dict[str, Any]] = []
    with pd.ExcelFile(path) as xls:
        for sheet_name in xls.sheet_names:
            df = xls.parse(sheet_name)
            md_table = _df_to_markdown(df, max_rows=max_rows)
            recs.append({"sheet": sheet_name, "source": "native", "section": "table", "text": md_table})
    if openpyxl is not None:
        try:
            wb = openpyxl.load_workbook(path, data_only=True)
            for ws in wb.worksheets:
                imgs = getattr(ws, "_images", [])
                for img in imgs:
                    try:
                        img_bytes = None
                        if hasattr(img, "_data") and callable(img._data):
                            img_bytes = img._data()
                        elif hasattr(img, "ref") and hasattr(img.ref, "blob"):
                            img_bytes = img.ref.blob
                        if not img_bytes:
                            continue
                        res = _analyze_image(img_bytes, ocr=ocr, vlm=vlm, prompt=None)
                        if res.get("ocr"):
                            recs.append({"sheet": ws.title, "source": "ocr", "section": "image", "text": res["ocr"]})
                        if res.get("caption"):
                            recs.append({"sheet": ws.title, "source": "vlm_alt", "section": "image_caption", "text": res["caption"]})
                    except Exception:
                        continue
                for row in ws.iter_rows(values_only=False):
                    for cell in row:
                        candidates: List[str] = []
                        try:
                            if cell.hyperlink is not None and getattr(cell.hyperlink, "target", None):
                                candidates.append(cell.hyperlink.target)
                        except Exception:
                            pass
                        v = cell.value
                        if isinstance(v, str):
                            candidates.append(v)
                        for cand in candidates:
                            b = _load_image_from_value(cand, image_root=image_root, fetch_http=fetch_http)
                            if not b:
                                continue
                            res = _analyze_image(b, ocr=ocr, vlm=vlm, prompt=None)
                            if res.get("ocr"):
                                recs.append({"sheet": ws.title, "source": "ocr", "section": "linked_image", "text": f"[cell {cell.coordinate}] {res['ocr']}"})
                            if res.get("caption"):
                                recs.append({"sheet": ws.title, "source": "vlm_alt", "section": "linked_image_caption", "text": f"[cell {cell.coordinate}] {res['caption']}"})
        except Exception:
            pass
    try:
        import zipfile
        with zipfile.ZipFile(path) as z:
            names = [n for n in z.namelist() if n.startswith("xl/media/")]
            for n in names:
                try:
                    b = z.read(n)
                    res = _analyze_image(b, ocr=ocr, vlm=vlm, prompt=None)
                    if res.get("ocr"):
                        recs.append({"sheet": "(unmapped)", "source": "ocr", "section": "image_unmapped", "text": res["ocr"]})
                    if res.get("caption"):
                        recs.append({"sheet": "(unmapped)", "source": "vlm_alt", "section": "image_unmapped_caption", "text": res["caption"]})
                except Exception:
                    continue
    except Exception:
        pass
    return recs

# ------------------------------------------------------
# YT-DLP DOWNLOAD → TRANSCRIBE (uses Step 3 ASR helpers at runtime)
# ------------------------------------------------------
def try_ytdlp_download_and_transcribe(url: str) -> Optional[str]:
    """
    Download video(s) via yt-dlp, extract audio with ffmpeg, transcribe with ASR,
    and return Markdown. If yt-dlp missing or nothing downloaded, return None.
    """
    try:
        import yt_dlp
    except Exception:
        return None  # yt-dlp not installed; let caller fallback

    with tempfile.TemporaryDirectory() as td:
        tdir = Path(td)
        ydl_opts = {
            "outtmpl": str(tdir / "%(title)s.%(ext)s"),
            "format": "best[is_dash=false][ext=mp4]/best[is_dash=false]/best",
            "noprogress": True,
            "quiet": True,
            "no_warnings": True,
        }
        try:
            with yt_dlp.YoutubeDL(ydl_opts) as ydl:
                ydl.download([url])
        except Exception:
            return None

        video_exts = {".mp4", ".mkv", ".webm", ".mov", ".m4v"}
        videos = [p for p in tdir.iterdir() if p.is_file() and p.suffix.lower() in video_exts]
        if not videos:
            return None

        md_sections: List[str] = []
        for vid in videos:
            wav_path = tdir / (vid.stem + ".wav")
            try:
                ffmpeg_convert_to_wav(vid, wav_path)
                segments, info = asr_transcribe_wav(wav_path)
                if not segments:
                    continue
                language = getattr(info, "language", "auto")
                md = segments_to_markdown(segments, vid.name, language)
                md_sections.append(md)
            finally:
                try:
                    if wav_path.exists(): wav_path.unlink()
                except Exception:
                    pass
                try:
                    if vid.exists(): vid.unlink()
                except Exception:
                    pass

        return "\n\n".join(md_sections) if md_sections else None

# ===========================================================================================================

# ======================================================
# STEP 3 — ASR PIPELINE (ffmpeg → Whisper → Markdown)
# ======================================================

def fmt_ts(sec: float) -> str:
    """Format seconds as HH:MM:SS.mmm (safe for None/NaN/inf)."""
    if sec is None or sec < 0 or math.isinf(sec) or math.isnan(sec):
        sec = 0.0
    ms = int(round((sec - int(sec)) * 1000))
    s  = int(sec) % 60
    m  = (int(sec) // 60) % 60
    h  = int(sec) // 3600
    return f"{h:02d}:{m:02d}:{s:02d}.{ms:03d}"

def ffmpeg_convert_to_wav(input_media: Path, output_wav: Path):
    """Use ffmpeg to turn any media file into mono 16kHz PCM WAV for ASR."""
    ffmpeg = get_ffmpeg_exe()
    cmd = [ffmpeg, "-y", "-i", str(input_media), "-vn", "-ac", "1", "-ar", "16000", "-acodec", "pcm_s16le", str(output_wav)]
    proc = subprocess.run(cmd, stdout=subprocess.PIPE, stderr=subprocess.PIPE)
    if proc.returncode != 0:
        raise RuntimeError(proc.stderr.decode("utf-8", "ignore"))

def asr_transcribe_wav(wav_path: Path):
    """Transcribe a WAV file with faster-whisper using VAD; return (segments, info)."""
    model = _get_asr_model()
    segments, info = model.transcribe(
        str(wav_path),
        vad_filter=True,
        vad_parameters=dict(min_silence_duration_ms=VAD_MIN_SIL_MS),
        beam_size=5,
        condition_on_previous_text=True,
    )
    return list(segments), info

def segments_to_markdown(segments, src_name: str, language: str) -> str:
    """Render Whisper segments into a simple, timestamped Markdown transcript."""
    header = (
        f"---\n"
        f"title: Transcript of {src_name}\n"
        f"generated: {datetime.datetime.now().isoformat(timespec='seconds')}\n"
        f"model: faster-whisper/{MODEL_SIZE} ({DEVICE},{COMPUTE_TYPE})\n"
        f"language: {language}\n"
        f"---\n\n# Audio Transcript\n\n"
    )
    body = []
    for seg in segments:
        body.append(f"{fmt_ts(seg.start)} --> {fmt_ts(seg.end)}\n{(seg.text or '').strip()}\n")
    return header + "\n".join(body)

def run_asr_pipeline_to_markdown(data: bytes, original_name: str) -> str:
    """
    Full ASR pipeline:
      1) write bytes to temp file,
      2) convert to WAV (ffmpeg),
      3) transcribe (VAD on; fallback pass if empty),
      4) lightly merge tiny gaps between segments,
      5) return Markdown report.
    """
    # Gap-merging knobs (seconds)
    GAP_S = float(os.getenv("ASR_MERGE_GAP_S", "0.5"))
    MAX_MERGED_UTTERANCE_S = float(os.getenv("ASR_MERGE_MAX_LEN_S", "60"))
    MIN_AUDIO_S = float(os.getenv("ASR_MIN_DURATION_S", "0.8"))

    class _Seg:
        __slots__ = ("start", "end", "text")
        def __init__(self, start: float, end: float, text: str):
            self.start = float(start)
            self.end   = float(end)
            self.text  = (text or "").strip()

    def _merge_segments(segs_in):
        segs = list(segs_in or [])
        if not segs:
            return []
        out = []
        cur = _Seg(segs[0].start, segs[0].end, segs[0].text)
        for s in segs[1:]:
            gap = max(0.0, float(s.start) - float(cur.end))
            merged_len = float(s.end) - float(cur.start)
            if gap <= GAP_S and merged_len <= MAX_MERGED_UTTERANCE_S:
                joiner = "" if (cur.text.endswith(("—", "-")) or cur.text.endswith(" ")) else " "
                cur.text = (cur.text + joiner + (s.text or "").strip()).strip()
                cur.end = max(cur.end, float(s.end))
            else:
                out.append(cur)
                cur = _Seg(s.start, s.end, s.text)
        out.append(cur)
        return out

    with tempfile.TemporaryDirectory() as td:
        td = Path(td)
        media_path = td / (original_name or "input.bin")
        wav_path   = td / "audio.wav"
        media_path.write_bytes(data)

        wav_created = False
        error_raised = False
        info = None
        segs = []  # ensure defined

        try:
            ffmpeg_convert_to_wav(media_path, wav_path)
            wav_created = wav_path.exists()

            # Pass 1 (default settings)
            segments, info = asr_transcribe_wav(wav_path)
            segs = list(segments)

            # Pass 2 (softer) if needed
            if not segs:
                segments2, info2 = ASR_MODEL.transcribe(
                    str(wav_path),
                    vad_filter=False,
                    beam_size=1,
                    temperature=0.7,
                    no_speech_threshold=0.2,
                    condition_on_previous_text=False,
                )
                segs = list(segments2)
                info = info2 or info

        except Exception:
            error_raised = True
            raise
        finally:
            # Keep WAV only if KEEP_WAV=true or (error and KEEP_WAV_ON_ERROR=true)
            keep_due_to_error = error_raised and os.getenv("KEEP_WAV_ON_ERROR", "true").lower() == "true"
            if wav_created and (not KEEP_WAV) and (not keep_due_to_error):
                try:
                    wav_path.unlink(missing_ok=True)
                except Exception:
                    pass
            # Remove original temp media bytes
            try:
                media_path.unlink(missing_ok=True)
            except Exception:
                pass

    # If nothing recognized
    duration = getattr(info, "duration", None)
    if not segs:
        language = getattr(info, "language", "auto") if info else "auto"
        note = "(Audio too short for transcription.)" if isinstance(duration, (int, float)) and duration < MIN_AUDIO_S \
               else "(No speech detected, or audio was too short/quiet.)"
        return (
            f"---\n"
            f"title: Transcript of {original_name}\n"
            f"generated: {datetime.datetime.now().isoformat(timespec='seconds')}\n"
            f"model: faster-whisper/{MODEL_SIZE} ({DEVICE},{COMPUTE_TYPE})\n"
            f"language: {language}\n"
            f"duration_s: {duration if duration is not None else 'unknown'}\n"
            f"---\n\n# Audio Transcript\n\n{note}\n"
        )

    # Merge tiny gaps for readability
    merged = _merge_segments(segs)
    language = getattr(info, "language", "auto")
    return segments_to_markdown(merged, original_name, language)


# ===========================================================================================================

# ======================================================
# STEP 4 — OCR (Florence-2 & PaddleOCR) + IMAGE→MARKDOWN
# ======================================================

def _parse_csv_env(name: str, default: str):
    """Read a comma-separated env var into a trimmed list."""
    raw = os.getenv(name, default)
    return [x.strip() for x in raw.split(",") if x.strip()]

# ------------------ User-configurable OCR knobs ------------------
OCR_LANGS = _parse_csv_env("OCR_LANGS", "en")
OCR_UPSCALE = float(os.getenv("OCR_UPSCALE", "1.5"))
OCR_RECOG_HEAD = os.getenv("OCR_RECOG_HEAD", "english_g2")  # compat only (unused here)
DEFAULT_ALLOWLIST = (
    "0123456789"
    "abcdefghijklmnopqrstuvwxyz"
    "ABCDEFGHIJKLMNOPQRSTUVWXYZ"
    " .,;:!?@#%&()[]{}+-/*=_'\"\\|<>~$^`"
)
OCR_ALLOWLIST = os.getenv("OCR_ALLOWLIST", DEFAULT_ALLOWLIST)

# ------------------ Small helpers ------------------
def _filter_allowlist(s: str) -> str:
    """Strip characters not in OCR_ALLOWLIST (keeps spaces)."""
    allowed = set(OCR_ALLOWLIST)
    return "".join(ch if (ch in allowed or ch.isspace()) else "" for ch in s)

def load_image_from_bytes(data: bytes) -> Image.Image:
    """Load bytes into RGB PIL Image; optionally upscale to help tiny text."""
    img = Image.open(io.BytesIO(data)).convert("RGB")
    if OCR_UPSCALE != 1.0:
        w, h = img.size
        img = img.resize((int(w * OCR_UPSCALE), int(h * OCR_UPSCALE)), Image.LANCZOS)
    return img

# ------------------ Florence-2 OCR execution ------------------
def _florence_ocr(img: Image.Image, task: str) -> Optional[dict]:
    """
    Run Florence-2 with a task token ('<OCR_WITH_REGION>' or '<OCR>').
    Returns Florence's parsed dict via processor.post_process_generation(...).
    """
    # Resolve model/processor (prefer lazy getter if available)
    try:
        model, proc = _get_florence()  # lazy path
    except Exception:
        model = globals().get("FLORENCE_MODEL", None)
        proc  = globals().get("FLORENCE_PROCESSOR", None)
    if model is None or proc is None:
        return None

    with torch.inference_mode():
        inputs = proc(text=task, images=img, return_tensors="pt")
        # Move tensors to device (dtype is already set on the model)
        inputs = {k: (v.to(FLORENCE_DEVICE) if hasattr(v, "to") else v) for k, v in inputs.items()}
        generated_ids = model.generate(
            input_ids=inputs["input_ids"],
            pixel_values=inputs["pixel_values"],
            max_new_tokens=1024,
            num_beams=3,
            do_sample=False,
            temperature=0.7,
            use_cache=True,
        )
        generated_text = proc.batch_decode(generated_ids, skip_special_tokens=False)[0]
        parsed = proc.post_process_generation(
            generated_text, task=task, image_size=(img.width, img.height)
        )
    return parsed


def _florence_extract_lines(parsed: dict) -> List[str]:
    """
    Normalize Florence parsed output into a list of text lines.
    Supports '<OCR_WITH_REGION>' (labels list) and '<OCR>' (single string).
    """
    if not isinstance(parsed, dict):
        return []
    # Prefer region OCR for better ordering when available
    for k, v in parsed.items():
        if "OCR_WITH_REGION" in k and isinstance(v, dict):
            labels = v.get("labels", []) or []
            return [t for t in (x.strip() for x in labels) if t]
    for k, v in parsed.items():
        if k.strip("<>") == "OCR" and isinstance(v, str):
            return [ln for ln in re.split(r"\s*[\r\n]+\s*", v.strip()) if ln]
    return []


def _florence_extract_regions(parsed: dict) -> List[Tuple[List[float], str]]:
    """
    From '<OCR_WITH_REGION>' parsed output, return [(bbox, label), ...].
    bbox is [x1, y1, x2, y2] in pixel coordinates when image_size was provided.
    """
    out: List[Tuple[List[float], str]] = []
    if not isinstance(parsed, dict):
        return out
    for k, v in parsed.items():
        if "OCR_WITH_REGION" in k and isinstance(v, dict):
            bboxes = v.get("bboxes", []) or []
            labels = v.get("labels", []) or []
            for bb, lb in zip(bboxes, labels):
                if isinstance(bb, (list, tuple)) and len(bb) == 4:
                    out.append((list(map(float, bb)), str(lb)))
    return out


def _crop_bbox(img: Image.Image, bbox: List[float]) -> Image.Image:
    """Crop image by [x1, y1, x2, y2] with clamping."""
    w, h = img.size
    x1 = max(0, min(w, int(round(bbox[0]))))
    y1 = max(0, min(h, int(round(bbox[1]))))
    x2 = max(0, min(w, int(round(bbox[2]))))
    y2 = max(0, min(h, int(round(bbox[3]))))
    if x2 <= x1 or y2 <= y1:
        return img.crop((0, 0, 1, 1))
    return img.crop((x1, y1, x2, y2))


# ------------------ Florence-2 caption execution ------------------
def _florence_caption(img: Image.Image, task: Optional[str] = None, max_new_tokens: Optional[int] = None) -> str:
    """
    Generate a semantic caption. Tries provided task (default: <MORE_DETAILED_CAPTION>),
    falls back to <DETAILED_CAPTION> then <CAPTION>. Returns '' on failure.
    """
    # Resolve model/processor (prefer lazy getter if available)
    try:
        model, proc = _get_florence()
    except Exception:
        model = globals().get("FLORENCE_MODEL", None)
        proc  = globals().get("FLORENCE_PROCESSOR", None)
    if model is None or proc is None:
        return ""

    task_primary = task or VLM_CAPTION_TASK
    max_tok = int(max_new_tokens or VLM_MAX_TOKENS_CAPTION)

    def _run(tkn: str) -> str:
        try:
            with torch.inference_mode():
                inputs = proc(text=tkn, images=img, return_tensors="pt")
                inputs = {k: (v.to(FLORENCE_DEVICE) if hasattr(v, "to") else v) for k, v in inputs.items()}
                ids = model.generate(
                    input_ids=inputs["input_ids"],
                    pixel_values=inputs["pixel_values"],
                    max_new_tokens=max_tok,
                    num_beams=3,
                    do_sample=False,
                    temperature=0.7,
                    use_cache=True,
                )
                text = proc.batch_decode(ids, skip_special_tokens=False)[0]
                parsed = proc.post_process_generation(
                    text, task=tkn, image_size=(img.width, img.height)
                )
            for k, v in (parsed or {}).items():
                key = k.strip("<>")
                if key in ("MORE_DETAILED_CAPTION", "DETAILED_CAPTION", "CAPTION") and isinstance(v, str):
                    return v.strip()
        except Exception:
            return ""
        return ""

    cap = _run(task_primary)
    if cap:
        return cap
    cap = _run("<DETAILED_CAPTION>")
    if cap:
        return cap
    return _run("<CAPTION>")

def _florence_region_captions(img: Image.Image, regions: List[Tuple[List[float], str]], topk: int) -> List[str]:
    """
    Caption up to top-k regions (by area). Returns list of strings:
    'Region i [x1,y1,x2,y2]: <caption>'
    """
    if not regions:
        return []
    w, h = img.size
    def _area(bb: List[float]) -> float:
        x1, y1, x2, y2 = bb
        return max(0.0, (min(x2, w) - max(x1, 0))) * max(0.0, (min(y2, h) - max(y1, 0)))
    regions_sorted = sorted(regions, key=lambda t: _area(t[0]), reverse=True)[:max(0, int(topk))]
    out: List[str] = []
    for i, (bb, _label) in enumerate(regions_sorted, start=1):
        try:
            crop = _crop_bbox(img, bb)
            rc = _florence_caption(crop, task=VLM_CAPTION_TASK, max_new_tokens=VLM_MAX_TOKENS_CAPTION)
            if rc:
                out.append(f"Region {i} [{int(bb[0])},{int(bb[1])},{int(bb[2])},{int(bb[3])}]: {rc}")
        except Exception:
            continue
    return out

# ------------------ PaddleOCR execution ------------------
def _paddle_ocr_lines(img: Image.Image) -> List[str]:
    """Run PaddleOCR (BGR input) and return recognized lines (allowlist-filtered)."""
    paddle_inst = _get_paddle()
    if not paddle_inst:
        return []
    img_np_bgr = np.array(img)[:, :, ::-1]
    result = paddle_inst.ocr(img_np_bgr, cls=True)
    lines: List[str] = []
    for per_image in (result or []):
        for line in (per_image or []):
            text = (line[1][0] if line and len(line) > 1 else "").strip()
            if text:
                text = _filter_allowlist(text)
                if text:
                    lines.append(text)
    return lines


# ------------------ Public: Image bytes → Markdown (Caption + OCR) ------------------
def ocr_image_bytes_to_markdown(data: bytes, src_name: str) -> str:
    """
    End-to-end for a single image:
      1) Global caption with Florence-2 (<MORE_DETAILED_CAPTION> → <DETAILED_CAPTION> → <CAPTION>),
      2) OCR with Florence (<OCR_WITH_REGION> → <OCR>) → PaddleOCR fallback,
      3) EXTRA: caption top-N detected regions for ultra-detail,
      4) Allowlist + de-dup,
      5) Render Markdown with caption, region captions, and text.
    """
    img = load_image_from_bytes(data)

    # 1) Global caption (extreme detail)
    caption = _florence_caption(img, VLM_CAPTION_TASK, VLM_MAX_TOKENS_CAPTION)
    cap_model = f"florence2/{VLM_MODEL_ID.split('/')[-1]}(more_detailed_caption)" if caption else "none"

    # 2) OCR (Florence → Paddle) + collect regions for region captions
    text_model = None
    lines: List[str] = []
    regions: List[Tuple[List[float], str]] = []
    parsed_region = None
    try:
        parsed_region = _florence_ocr(img, "<OCR_WITH_REGION>")
        lines = _florence_extract_lines(parsed_region)
        regions = _florence_extract_regions(parsed_region)
        text_model = f"florence2/{VLM_MODEL_ID.split('/')[-1]}(ocr_region)"
        if not any(ch.isalnum() for ln in lines for ch in ln):
            parsed2 = _florence_ocr(img, "<OCR>")
            lines2 = _florence_extract_lines(parsed2)
            if lines2:
                lines = lines2
                text_model = f"florence2/{VLM_MODEL_ID.split('/')[-1]}(ocr)"
    except Exception:
        lines, regions = [], []

    if not lines or not any(ch.isalnum() for ln in lines for ch in ln):
        try:
            lines = _paddle_ocr_lines(img)
            text_model = f"paddleocr/{PADDLE_LANG}"
        except Exception:
            lines, text_model = [], "none"

    # 3) Region captions (top-N by area)
    region_caps = _florence_region_captions(img, regions, VLM_REGION_CAPTION_TOPK) if regions else []

    # 4) Allowlist + simple de-dup
    filtered = []
    seen = set()
    for ln in lines:
        f = _filter_allowlist(ln).strip()
        if f and f not in seen:
            seen.add(f)
            filtered.append(f)

    # 5) Render Markdown
    header = (
        f"---\n"
        f"title: OCR/Caption of {src_name}\n"
        f"generated: {datetime.datetime.now().isoformat(timespec='seconds')}\n"
        f"caption_model: {cap_model}\n"
        f"text_model: {text_model} "
        f"(vlm_device={'cuda' if torch.cuda.is_available() else 'cpu'}, "
        f"paddle_device={'cuda' if PADDLE_USE_GPU else 'cpu'})\n"
        f"vlm_caption_task: {VLM_CAPTION_TASK}\n"
        f"vlm_caption_tokens: {VLM_MAX_TOKENS_CAPTION}\n"
        f"region_captions: {len(region_caps)} (top_k={VLM_REGION_CAPTION_TOPK})\n"
        f"language: {','.join(OCR_LANGS)}\n"
        f"---\n\n"
    )

    parts = []
    if caption:
        parts.append("# Image Caption\n\n" + caption.strip())
    if region_caps:
        parts.append("## Region Captions\n\n" + "\n".join(f"- {rc}" for rc in region_caps))
    parts.append("# OCR Text\n\n" + ("\n\n".join(filtered) if filtered else "(no text found)"))
    return header + "\n\n".join(parts) + "\n"

# ===========================================================================================================

# ======================================================
# STEP 5 — HTML (meta / clean / single-page) + URL fetch
# ======================================================

def _meta(soup: BeautifulSoup, key: str) -> Optional[str]:
    """Return a meta value by name/og:/twitter: precedence, or None if missing."""
    tag = (
        soup.find("meta", attrs={"name": key})
        or soup.find("meta", attrs={"property": f"og:{key}"})
        or soup.find("meta", attrs={"property": f"twitter:{key}"})
    )
    return tag.get("content", "").strip() if tag and tag.has_attr("content") else None

def _clean(
    text: str,
    *,
    normalize_unicode: bool = True,
    collapse_whitespace: bool = True,
    strip_line_ends: bool = True,
    dedent: bool = True,
) -> str:
    """
    Normalize raw page text for Markdown:
      - normalize Unicode (NFKC), drop zero-widths
      - unify newlines and strip trailing spaces
      - dedent code-like blocks
      - collapse runs of spaces and blank lines
    """
    if normalize_unicode:
        text = unicodedata.normalize("NFKC", text)
        text = text.replace("\u00A0", " ")
        text = re.sub(r"[\u200B-\u200D\u2060\uFEFF]", "", text)
    text = text.replace("\r\n", "\n").replace("\r", "\n")
    if strip_line_ends:
        text = "\n".join(line.rstrip() for line in text.splitlines())
    if dedent:
        text = textwrap.dedent(text)
    if collapse_whitespace:
        text = re.sub(r"[ \t]+", " ", text)
        text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()

def single_page_markdown(url: str, html_bytes: bytes) -> str:
    """
    Convert a single HTML page to Markdown:
      - parse <title>, author/description meta
      - extract visible text and clean it
      - return a Markdown doc with front matter
    """
    soup = BeautifulSoup(html_bytes, "lxml")
    title = (soup.title.string.strip() if soup.title else url)
    description = _meta(soup, "description")
    author = _meta(soup, "author")
    content = _clean(soup.get_text(" ", strip=True))

    header = (
        f"---\n"
        f"title: Page extract of {title}\n"
        f"generated: {datetime.datetime.now().isoformat(timespec='seconds')}\n"
        f"url: {url}\n"
        f"pages: 1\n"
        f"---\n\n# Page Content\n\n"
    )
    meta_block = []
    if author: meta_block.append(f"**Author:** {author}")
    if description: meta_block.append(f"**Description:** {description}")
    meta_str = ("\n\n" + "\n\n".join(meta_block) + "\n\n") if meta_block else "\n"
    return header + meta_str + content + "\n"

async def fetch_head_or_get(url: str, user_agent: Optional[str] = None) -> Tuple[bytes, str]:
    """
    Fetch a URL with a HEAD probe (to see Content-Type) and then GET the body.
    Returns (content_bytes, lowercase_content_type).

    Behavior:
    - Adds https:// if scheme missing (only http/https allowed).
    - Enforces per-request timeout via env HTTP_TIMEOUT_S (default 30s).
    - Enforces max download size via env MAX_FETCH_BYTES (default 25 MiB); raises 413 if exceeded.
    - On DNS/connect/timeout issues, raises HTTP 502 with a friendly message.
    - On non-200 GET, raises HTTP 400.
    """
    MAX_FETCH_BYTES = int(os.getenv("MAX_FETCH_BYTES", str(25 * 1024 * 1024)))  # 25 MiB
    HTTP_TIMEOUT_S = float(os.getenv("HTTP_TIMEOUT_S", "30"))

    headers = {"User-Agent": user_agent} if user_agent else {}

    # Normalize URL and validate scheme
    if "://" not in url:
        url = "https://" + url
    try:
        scheme = httpx.URL(url).scheme
    except Exception:
        raise HTTPException(400, "Invalid URL")
    if scheme not in ("http", "https"):
        raise HTTPException(400, "Unsupported URL scheme; only http(s) is allowed")

    timeout = httpx.Timeout(HTTP_TIMEOUT_S)

    async with httpx.AsyncClient(follow_redirects=True, timeout=timeout, headers=headers) as client:
        try:
            ctype = ""

            # Try HEAD first (best-effort; ignore failures and continue to GET)
            try:
                hr = await client.head(url)
                ctype = (hr.headers.get("Content-Type") or "").lower()
                clen_hdr = hr.headers.get("Content-Length")
                if clen_hdr:
                    try:
                        clen = int(clen_hdr)
                        if clen > MAX_FETCH_BYTES:
                            raise HTTPException(413, f"Remote file too large ({clen} bytes > {MAX_FETCH_BYTES}).")
                    except Exception:
                        pass  # bad/missing length → fall through
            except Exception:
                pass  # some servers don't support HEAD; proceed to GET

            # Stream GET and enforce size cap
            async with client.stream("GET", url) as gr:
                if gr.status_code != 200:
                    raise HTTPException(400, f"Failed to fetch URL (status {gr.status_code})")

                ctype = (gr.headers.get("Content-Type") or ctype).lower()

                clen_hdr = gr.headers.get("Content-Length")
                if clen_hdr:
                    try:
                        clen = int(clen_hdr)
                        if clen > MAX_FETCH_BYTES:
                            raise HTTPException(413, f"Remote file too large ({clen} bytes > {MAX_FETCH_BYTES}).")
                    except Exception:
                        pass

                buf = bytearray()
                async for chunk in gr.aiter_bytes():
                    if not chunk:
                        continue
                    buf.extend(chunk)
                    if len(buf) > MAX_FETCH_BYTES:
                        raise HTTPException(413, f"Downloaded file too large (> {MAX_FETCH_BYTES} bytes).")

                return bytes(buf), ctype

        except ConnectError as e:
            # Clear message for DNS / connect errors
            raise HTTPException(502, f"Network error while fetching URL (DNS/connect): {e}") from e
        except httpx.TimeoutException as e:
            raise HTTPException(502, f"Network timeout while fetching URL: {e}") from e
        except HTTPException:
            raise  # bubble up our own HTTP errors unchanged
        except Exception as e:
            raise HTTPException(502, f"Network error while fetching URL: {e}") from e

# ===========================================================================================================

# ======================================================
# STEP 6 — PUBLIC DRIVER (mdify_ingest_process_to_markdown)
# ======================================================

class _IngestArgs:
    """Lightweight args holder: access kwargs as attributes (e.g., args.foo)."""
    def __init__(self, **kw):
        for k, v in kw.items():
            setattr(self, k, v)


def mdify_ingest_process_to_markdown(path: str, *,
    outdir: Optional[str] = None,
    title: Optional[str] = None,
    title_prefix: Optional[str] = None,
    url: Optional[str] = None,
    encoding: Optional[str] = None,
    ocr: str = "paddle",
    lang: str = "en",
    dpi: int = 300,
    ocr_threshold: int = 50,
    vlm: str = "hf_local",
    vlm_model: Optional[str] = None,
    vlm_max_tokens: int = 128,
    vlm_prompt: str = "<MORE_DETAILED_CAPTION>",  # ← default to more-detailed captions
    caption_pages: bool = False,
    vlm_dpi: int = 256,
    image_columns: Optional[str] = None,
    fetch_http: bool = False,
    image_root: str = ".",
    caption_xlsx_charts: bool = False,
    soffice_bin: Optional[str] = None,
    max_rows: int = -1,  # <=0 means "all rows"
    vlm_only: bool = False,  # ← caption-only mode for pdf/pptx/xlsx/doc/docx
) -> str:
    """
    Ingest a single file and emit a Markdown note with front matter + grouped contents.
    """
    p = pathlib.Path(path)
    ext = p.suffix.lower()
    if ext not in _DOC_TYPES:
        raise ValueError(f"Unsupported file type: {ext}")
    mime, doc_type = _DOC_TYPES[ext]

    # Normalize row-limit: <=0 → all rows (propagate as -1 to extractors/_df_to_markdown)
    row_limit = -1 if (max_rows is None or int(max_rows) <= 0) else int(max_rows)
    row_limit_meta = "all" if row_limit <= 0 else row_limit  # for front-matter only

    # Engines
    ocr_wrap = _OCRWrapper(engine=(None if ocr == "none" else "paddle"), lang=lang)

    # Default the Florence model if VLM enabled
    if vlm == "hf_local" and not vlm_model:
        vlm_model = "microsoft/Florence-2-large"

    # Make sure the wrapper uses our (possibly more-detailed) prompt
    vlm_wrap = _VLMWrapper(
        model=vlm_model if vlm == "hf_local" or vlm_only else None,
        max_tokens=vlm_max_tokens,
        prompt=vlm_prompt
    )

    # ---------- helpers ----------
    def _caption_pages_via_pdf(pdf_path: str) -> List[Dict[str, Any]]:
        if vlm_wrap and vlm_wrap.enabled:
            try:
                return _caption_pdf_pages(pdf_path, vlm_wrap, dpi=vlm_dpi)
            except Exception:
                return []
        return []

    def _export_then_caption(export_fn_name: str, src_path: str) -> List[Dict[str, Any]]:
        """Call a soffice exporter if available, then caption pages of the produced PDF."""
        export_fn = globals().get(export_fn_name)
        if callable(export_fn):
            try:
                pdf_path = export_fn(src_path, soffice_bin=soffice_bin)
                if pdf_path and os.path.exists(pdf_path):
                    return _caption_pages_via_pdf(pdf_path)
            except Exception:
                pass
        return []

    def _filter_vlm_only(recs: List[Dict[str, Any]]) -> List[Dict[str, Any]]:
        return [r for r in (recs or []) if r.get("source") == "vlm_alt"]

    # ---------- caption-only fast paths ----------
    if vlm_only:
        if not (vlm_wrap and vlm_wrap.enabled):
            raise RuntimeError("vlm_only=True requires a local VLM; set vlm='hf_local' and provide/allow a model.")
        if ext == ".pdf":
            records = _caption_pages_via_pdf(path)
        elif ext == ".pptx":
            records = _export_then_caption("_export_pptx_to_pdf", path)
            if not records:
                recs = _extract_pptx(path, ocr=_OCRWrapper(engine=None, lang=lang), vlm=vlm_wrap)
                records = _filter_vlm_only(recs)
        elif ext == ".docx":
            records = _export_then_caption("_export_docx_to_pdf", path)
            if not records:
                recs = _extract_docx(path, ocr=_OCRWrapper(engine=None, lang=lang), vlm=vlm_wrap)
                records = _filter_vlm_only(recs)
        elif ext == ".doc":
            records = _export_then_caption("_export_doc_to_pdf", path)
            if not records:
                recs = _extract_doc_legacy(path, ocr=_OCRWrapper(engine=None, lang=lang), vlm=vlm_wrap)
                records = _filter_vlm_only(recs)
        elif ext == ".xlsx":
            records = []
            try:
                pdf_path = _export_xlsx_to_pdf(path, soffice_bin=soffice_bin)
                if pdf_path:
                    records = _caption_pages_via_pdf(pdf_path)
            except Exception:
                records = []
            if not records:
                recs = _extract_xlsx_enhanced(path, ocr=_OCRWrapper(engine=None, lang=lang), vlm=vlm_wrap)
                records = _filter_vlm_only(recs)
        elif ext in (".txt", ".csv"):
            records = _extract_txt(path, encoding=encoding) if ext == ".txt" else _extract_csv_enhanced(
                path, ocr=None, vlm=vlm_wrap, encoding=encoding, max_rows=row_limit,
                image_columns=image_columns, image_root=image_root, fetch_http=fetch_http
            )
            records = _filter_vlm_only(records)
        else:
            raise AssertionError("unreachable")

    # ---------- full (normal) pipeline ----------
    else:
        if ext == ".pdf":
            records = _extract_pdf(
                path,
                ocr=ocr_wrap,
                vlm=vlm_wrap,
                dpi=dpi,
                char_threshold=ocr_threshold,
                caption_pages=caption_pages,
                vlm_dpi=vlm_dpi,
                region_dpi=384,
                image_limit_per_page=256,
                min_region_wh=22,
                small_region_dpi=640,
                small_region_pad=2,
            )
            if vlm_wrap and vlm_wrap.enabled and caption_pages:
                try:
                    records.extend(_caption_pdf_pages(path, vlm_wrap, dpi=vlm_dpi))
                except Exception:
                    pass

        elif ext == ".pptx":
            records = _extract_pptx(
                path,
                ocr=ocr_wrap,
                vlm=vlm_wrap,
                tiny_min_px=22,
                target_min_px=64,
                max_icons_per_slide=512,
            )
            if caption_pages and vlm_wrap and vlm_wrap.enabled:
                records.extend(_export_then_caption("_export_pptx_to_pdf", path))

        elif ext == ".doc":
            records = _extract_doc_legacy(
                path,
                ocr=ocr_wrap,
                vlm=vlm_wrap,
                soffice_bin=soffice_bin,
                dpi=dpi,
                ocr_threshold=ocr_threshold,
                caption_pages=caption_pages,
                vlm_dpi=vlm_dpi,
            )
            if caption_pages and vlm_wrap and vlm_wrap.enabled:
                records.extend(_export_then_caption("_export_doc_to_pdf", path))

        elif ext == ".docx":
            records = _extract_docx(
                path,
                ocr=ocr_wrap,
                vlm=vlm_wrap,
                tiny_min_px=22,
                target_min_px=64,
                max_images=1024,
            )
            if caption_pages and vlm_wrap and vlm_wrap.enabled:
                records.extend(_export_then_caption("_export_docx_to_pdf", path))

        elif ext == ".txt":
            records = _extract_txt(path, encoding=encoding)

        elif ext == ".csv":
            records = _extract_csv_enhanced(
                path,
                ocr=ocr_wrap,
                vlm=vlm_wrap,
                encoding=encoding,
                max_rows=row_limit,
                image_columns=image_columns,
                image_root=image_root,
                fetch_http=fetch_http,
            )

        elif ext == ".xlsx":
            records = _extract_xlsx_enhanced(
                path,
                ocr=ocr_wrap,
                vlm=vlm_wrap,
                max_rows=row_limit,
                image_root=image_root,
                fetch_http=fetch_http,
            )
            if (vlm_wrap and vlm_wrap.enabled) and (caption_xlsx_charts or caption_pages):
                try:
                    pdf_path = _export_xlsx_to_pdf(path, soffice_bin=soffice_bin)
                    if pdf_path:
                        records.extend(_caption_pages_via_pdf(pdf_path))
                except Exception:
                    pass
        else:
            raise AssertionError("unreachable")

    # Title & render
    title_final = title or f"Ingest of {p.name}"
    if title_prefix:
        title_final = f"{title_prefix}: {title_final}"

    times_meta = _collect_doc_times(path, doc_type)

    note = _render_markdown_note(
        title=title_final,
        src_path=path,
        url=url,
        records=records,
        mime=mime,
        doc_type=doc_type,
        languages=[lang],
        text_extractor=Path(__file__).name,   # ← reflect current module
        ocr_engine=("paddle" if ocr_wrap.engine else ""),
        captioner=(vlm_model or ""),
        extra_meta={
            "row_limit": row_limit_meta,     # ← emit the normalized meta value
            "timestamps": times_meta,
            "vlm_active": bool(vlm_wrap and vlm_wrap.enabled),
        },
    )
    return note

# ===========================================================================================================

# ======================================================
# STEP 7 — FASTAPI APP (exception handler, models, routes)
# ======================================================

app = FastAPI(
    title="Unified ASR/OCR & Document Ingest Server",
    version="0.1",
    docs_url="/__docs__",
    openapi_url="/openapi.json",   # explicit is nice for proxies
)

# NEW – permissive for local dev; tighten in prod
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],           # e.g. ["http://localhost:3000"]
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# ---- Preload models before the app starts serving ----
@app.on_event("startup")
def _preload_models():
    import time
    if os.getenv("PRELOAD_MODELS", "true").lower() != "true":
        print("[startup] PRELOAD_MODELS=false → skipping preloads.", flush=True)
        return

    # 1) Whisper first
    t0 = time.time()
    try:
        _ = _get_asr_model()  # instantiates WhisperModel once
        dt = time.time() - t0
        print(f"[startup] Whisper ready ({DEVICE},{COMPUTE_TYPE}) in {dt:.1f}s", flush=True)
    except Exception as e:
        print(f"[startup] Whisper preload FAILED: {e}", file=sys.stderr, flush=True)

    # 2) Florence next (snapshot + model+processor)
    t1 = time.time()
    try:
        try:
            _ = _ensure_florence_path()  # make sure local snapshot exists
        except Exception:
            pass

        m, p = _get_florence()
        dt = time.time() - t1
        if m and p:
            dtype = str(FLORENCE_DTYPE).replace("torch.", "")
            print(f"[startup] Florence ready ({FLORENCE_DEVICE},{dtype}) in {dt:.1f}s", flush=True)
        else:
            print(f"[startup] Florence preload SKIPPED — last_error={_FLORENCE_LAST_ERROR}", file=sys.stderr, flush=True)
    except Exception as e:
        print(f"[startup] Florence preload FAILED: {e}", file=sys.stderr, flush=True)


# ---------- Friendly error handler ----------
@app.exception_handler(HTTPException)
async def http_exception_handler(request: Request, exc: HTTPException):
    """Return a plain 502 for DNS/connect issues; JSON for other HTTP errors."""
    detail = str(exc.detail or "")
    if exc.status_code == 502 and "Network error while fetching URL" in detail:
        return PlainTextResponse(
            "We couldn’t reach that URL. Please confirm it’s typed correctly, accessible from your network, and try again.",
            status_code=502,
            media_type="text/plain",
        )
    return JSONResponse(status_code=exc.status_code, content={"detail": exc.detail})

# ---------- Models ----------
class UrlBody(BaseModel):
    """Simple body schema for URL-based transcription/extraction."""
    url: str
    user_agent: Optional[str] = None  # optional UA for the single GET

# Default UA for outbound fetches (overridable via env FETCH_UA)
DEFAULT_FETCH_UA = os.getenv(
    "FETCH_UA",
    "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/120 Safari/537.36",
)

# ---- URL/MIME → supported doc extension helpers ----
_DOC_EXT_BY_MIME = {mime: ext for ext, (mime, _) in _DOC_TYPES.items()}

def _pick_supported_ext_from_url_or_ctype(url: str, ctype: str) -> Optional[str]:
    """
    Return a supported extension ('.pdf', '.docx', '.pptx', '.xlsx', '.csv') if
    the URL path or HTTP Content-Type matches a known document type; else None.
    """
    # Prefer by URL suffix
    try:
        parsed = httpx.URL(url if "://" in url else "https://" + url)
        ext = Path(parsed.path).suffix.lower()
        if ext in _DOC_TYPES:
            return ext
    except Exception:
        pass
    # Fallback by MIME (strip charset)
    mime = (ctype or "").split(";")[0].lower().strip()
    return _DOC_EXT_BY_MIME.get(mime)

# ================== Routes ===================
@app.get("/health", response_class=PlainTextResponse)
def health():
    """Liveness probe."""
    return "ok"

@app.get("/", response_class=PlainTextResponse)
def root():
    return (
        f"ASR: faster-whisper/{MODEL_SIZE} ({DEVICE},{COMPUTE_TYPE}) | "
        f"VLM OCR+Caption: Florence-2 ({'cuda' if torch.cuda.is_available() else 'cpu'}) + "
        f"PaddleOCR({ 'cuda' if PADDLE_USE_GPU else 'cpu'}) "
        f"[langs={','.join(OCR_LANGS)}] | "
        f"URL: single-page extract + yt-dlp media"
    )

@app.get("/__vlm_status")
def __vlm_status():
    try:
        m, p = _get_florence()
        ok = bool(m and p)
    except Exception as e:
        ok = False
    return JSONResponse({
        "enabled": ok,
        "model_id": VLM_MODEL_ID,
        "local_dir": FLORENCE_LOCAL_DIR,
        "device": FLORENCE_DEVICE,
        "dtype": str(FLORENCE_DTYPE),
        "hf_home": os.getenv("HF_HOME") or os.getenv("TRANSFORMERS_CACHE"),
        "offline": {
            "TRANSFORMERS_OFFLINE": os.getenv("TRANSFORMERS_OFFLINE"),
            "HF_HUB_OFFLINE": os.getenv("HF_HUB_OFFLINE"),
        },
        "last_error": _FLORENCE_LAST_ERROR,
    })

@app.get("/__asr_status")
def __asr_status():
    return JSONResponse({
        "model_size": MODEL_SIZE,
        "device": DEVICE,
        "compute_type": COMPUTE_TYPE,
        "snapshot": {
            "enabled": WHISPER_SNAPSHOT_ENABLED,
            "dir": str(WHISPER_SNAPSHOT_DIR),
            "allow_download": WHISPER_ALLOW_DOWNLOAD,
            "revision": WHISPER_REVISION,
            "last_error": _WHISPER_SNAPSHOT_LAST_ERROR,
        },
    })


@app.post("/extract_document")
async def extract_document(
    request: Request,
    # Metadata
    title: Optional[str] = None,
    title_prefix: Optional[str] = None,
    url: Optional[str] = None,
    encoding: Optional[str] = None,
    # OCR
    ocr: str = "paddle",  # "paddle" or "none"
    lang: str = "en",
    dpi: int = 300,
    ocr_threshold: int = 50,
    # VLM
    vlm: str = "hf_local",  # "hf_local" or "none"
    vlm_model: Optional[str] = None,  # default microsoft/Florence-2-large
    vlm_max_tokens: int = 128,
    vlm_prompt: str = "<MORE_DETAILED_CAPTION>",  # ← align with core default
    caption_pages: bool = False,
    vlm_dpi: int = 256,
    # CSV/XLSX images
    image_columns: Optional[str] = None,
    fetch_http: bool = False,
    image_root: str = ".",
    # XLSX charts via LibreOffice
    caption_xlsx_charts: bool = False,
    soffice_bin: Optional[str] = None,
    # Tables
    max_rows: int = -1,  # <=0 means "ALL rows" (no truncation)
    # Caption-only fast path
    vlm_only: bool = False,  # ← expose caption-only mode
):
    """Accept a binary doc and return a Markdown extraction (text + OCR + captions)."""
    data = await request.body()
    if not data:
        raise HTTPException(400, "Empty request body")

    name = request.headers.get("X-Filename", "document.bin")
    ext = os.path.splitext(name)[1].lower()
    if ext not in _DOC_TYPES:
        raise HTTPException(
            415,
            f"Unsupported file type: {ext or 'unknown'}; expected one of {sorted(_DOC_TYPES.keys())}"
        )

    with tempfile.TemporaryDirectory() as td:
        td_path = Path(td)
        in_path = td_path / name
        in_path.write_bytes(data)
        try:
            md = mdify_ingest_process_to_markdown(
                str(in_path),
                outdir=str(td_path),
                title=title,
                title_prefix=title_prefix,
                url=url,
                encoding=encoding,
                ocr=ocr,
                lang=lang,
                dpi=dpi,
                ocr_threshold=ocr_threshold,
                vlm=vlm,
                vlm_model=vlm_model,
                vlm_max_tokens=vlm_max_tokens,
                vlm_prompt=vlm_prompt,
                caption_pages=caption_pages,
                vlm_dpi=vlm_dpi,
                image_columns=image_columns,
                fetch_http=fetch_http,
                image_root=image_root,
                caption_xlsx_charts=caption_xlsx_charts,
                soffice_bin=soffice_bin,
                max_rows=max_rows,    # pass through (<=0 → ALL handled downstream)
                vlm_only=vlm_only,    # ← pass through
            )
        except Exception as e:
            raise HTTPException(500, f"Ingest failed: {e}") from e

    return PlainTextResponse(md, media_type="text/markdown")


@app.post("/transcribe_file")
async def transcribe_file(request: Request):
    """Accept raw media bytes, run ASR, return Markdown transcript."""
    data = await request.body()
    if not data:
        raise HTTPException(400, "Empty request body")
    name = request.headers.get("X-Filename", "input.bin")
    md = run_asr_pipeline_to_markdown(data, name)
    return PlainTextResponse(md, media_type="text/markdown")

@app.post("/transcribe_image")
async def transcribe_image(request: Request):
    """Accept raw image bytes, run OCR (Florence→Paddle fallback), return Markdown."""
    data = await request.body()
    if not data:
        raise HTTPException(400, "Empty request body")
    name = request.headers.get("X-Filename", "image.bin")
    md = ocr_image_bytes_to_markdown(data, name)
    return PlainTextResponse(md, media_type="text/markdown")

@app.post("/transcribe_url")
async def transcribe_url(body: UrlBody):
    """
    URL workflow:
      1) Try yt-dlp for video pages (download→ffmpeg→ASR).
      2) Else single GET:
         - Documents (PDF/DOCX/PPTX/XLSX/CSV) → document ingester
         - HTML → single-page markdown
         - image/* or image by sniff → OCR
         - audio/video/* or media by sniff → ASR
    """
    # ------- tiny helpers (scoped) -------
    def _sniff_ext_from_bytes(b: bytes) -> Optional[str]:
        # PDF
        if b.startswith(b"%PDF-"):
            return ".pdf"
        # OOXML (zip) → check [Content_Types].xml
        if b[:2] == b"PK":
            try:
                import zipfile, io as _io
                with zipfile.ZipFile(_io.BytesIO(b)) as z:
                    ct = z.read("[Content_Types].xml").decode("utf-8", "ignore")
                    if "wordprocessingml.document" in ct: return ".docx"
                    if "presentationml.presentation" in ct: return ".pptx"
                    if "spreadsheetml.sheet" in ct: return ".xlsx"
            except Exception:
                pass
        # Basic image headers
        if b[:8] == b"\x89PNG\r\n\x1a\n": return ".png"
        if b[:3] == b"\xff\xd8\xff": return ".jpg"
        if b[:6] == b"GIF87a" or b[:6] == b"GIF89a": return ".gif"
        return None

    def _looks_like_media(b: bytes) -> bool:
        # Audio: WAV/RIFF, MP3 (ID3), OGG/FLAC; Video: MP4/ISO BMFF, WebM/Matroska
        if b.startswith(b"RIFF") and b[8:12] in (b"WAVE", b"AVI "):
            return True
        if b.startswith(b"ID3"):  # MP3 (ID3 tag)
            return True
        if b.startswith(b"OggS") or b.startswith(b"fLaC"):
            return True
        if len(b) >= 12 and b[4:8] == b"ftyp":  # MP4/ISO BMFF
            return True
        if b.startswith(b"\x1A\x45\xDF\xA3"):  # Matroska/WebM
            return True
        return False
    # -------------------------------------

    url = (body.url or "").strip()
    if not url:
        raise HTTPException(400, "Missing 'url'")

    # 1) Prefer yt-dlp for video pages
    md_from_ytdlp = try_ytdlp_download_and_transcribe(url)
    if md_from_ytdlp:
        return PlainTextResponse(md_from_ytdlp, media_type="text/markdown")

    # 2) Single fetch (clear 502 on DNS/connect handled upstream)
    data, ctype = await fetch_head_or_get(url, body.user_agent or DEFAULT_FETCH_UA)
    if not data:
        raise HTTPException(400, "Fetched URL returned empty body")

    MAX_FETCH_BYTES = int(os.getenv("MAX_FETCH_BYTES", str(25 * 1024 * 1024)))
    if len(data) > MAX_FETCH_BYTES:
        raise HTTPException(413, f"Downloaded file too large ({len(data)} bytes > {MAX_FETCH_BYTES}).")

    url_norm = url if "://" in url else "https://" + url

    # ---- Decide by type/sniff (docs vs images vs media) ----
    sniff_ext = _sniff_ext_from_bytes(data)

    # 2a) Supported office/web docs → mdify_ingest_process_to_markdown
    DOC_ONLY = {".pdf", ".docx", ".pptx", ".xlsx", ".csv"}  # NOTE: no images here
    doc_ext = _pick_supported_ext_from_url_or_ctype(url, (ctype or "")) or (sniff_ext if sniff_ext in DOC_ONLY else None)
    if doc_ext:
        with tempfile.TemporaryDirectory() as td:
            parsed = httpx.URL(url_norm)
            name = Path(parsed.path).name or f"remote{doc_ext}"
            if not name.lower().endswith(doc_ext):
                name = f"{name}{doc_ext}"
            in_path = Path(td) / name
            in_path.write_bytes(data)
            md = mdify_ingest_process_to_markdown(
                str(in_path),
                url=url,
                ocr="paddle",
                lang="en",
            )
        return PlainTextResponse(md, media_type="text/markdown")

    # 2b) HTML page → single-page extract
    if is_probably_html((ctype or ""), data):
        md = single_page_markdown(url_norm, data)
        return PlainTextResponse(md, media_type="text/markdown")

    # 2c) Image → OCR (by Content-Type OR by sniff)
    is_img_ctype = bool(ctype) and is_image_content_type(ctype)
    is_img_sniff = sniff_ext in {".png", ".jpg", ".gif"}
    if is_img_ctype or is_img_sniff:
        try:
            Image.open(io.BytesIO(data)).verify()  # confirm it's an image
            name = Path(httpx.URL(url_norm).path).name or "remote_image"
            md = ocr_image_bytes_to_markdown(data, name)
            return PlainTextResponse(md, media_type="text/markdown")
        except Exception:
            pass

    # 2d) Audio/Video → ASR strictly on media signal (no generic octet-stream)
    if (ctype and is_audio_or_video_content_type(ctype) and ctype.lower() != "application/octet-stream") or _looks_like_media(data):
        name = Path(httpx.URL(url_norm).path).name or "remote_media"
        md = run_asr_pipeline_to_markdown(data, name)
        return PlainTextResponse(md, media_type="text/markdown")

    # Unknown
    raise HTTPException(415, f"Unsupported Content-Type for URL: {ctype or 'unknown'}")
